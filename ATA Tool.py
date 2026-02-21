import io
import base64
import json
import os
import sys
import tempfile
from datetime import date, datetime
from uuid import uuid4
from pathlib import Path
from urllib.parse import quote
import matplotlib.pyplot as plt
import pandas as pd
import gspread
import streamlit as st
from google.oauth2.service_account import Credentials
import streamlit.components.v1 as components
import extra_streamlit_components as stx
from fpdf import FPDF
from openpyxl import load_workbook
from openpyxl.styles import Alignment, Font, PatternFill
from pptx import Presentation
from pptx.util import Inches
# -------------------- APP CONFIG --------------------
st.set_page_config(page_title="DAMAC | ATA Tool", layout="wide")


def secure_authentication_gate() -> None:
    """Secure login gate using Streamlit secrets with lockout/cooldown."""
    if "authenticated" not in st.session_state:
        st.session_state.authenticated = False
    if "login_attempts" not in st.session_state:
        st.session_state.login_attempts = 0
    if "lockout_until" not in st.session_state:
        st.session_state.lockout_until = None

    if st.session_state.authenticated:
        return

    expected_username = st.secrets["auth"]["username"]
    expected_password = st.secrets["auth"]["password"]

    now_utc = datetime.utcnow()
    lockout_until = st.session_state.lockout_until
    locked = bool(lockout_until and now_utc < lockout_until)
    remaining = int((lockout_until - now_utc).total_seconds()) if locked else 0

    st.markdown(
        f"""
        <div class="login-form-shell">
            <div class="login-logo-card">
                <div class="login-logo"><img src="{LOGIN_LOGO_URL}" style="max-width:90%;height:auto;"></div>
                <h3 class="login-title">Welcome To ATA Tool</h3>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    left, center, right = st.columns([1, 2, 1])
    with center:
        st.markdown('<div class="login-form-shell">', unsafe_allow_html=True)
        with st.form("secure_login_form"):
            username = st.text_input("User Name", placeholder="Enter username")
            password = st.text_input("Password", type="password", placeholder="Enter password")
            submitted = st.form_submit_button("Login", use_container_width=True, disabled=locked)
        st.markdown('</div>', unsafe_allow_html=True)

    if locked:
        st.error(f"Too many failed attempts. Please wait {remaining} seconds before trying again.")

    st.markdown(
        "<div class='login-extra'>This App was created for quality activity purposes.</div>",
        unsafe_allow_html=True,
    )

    if submitted and not locked:
        if username == expected_username and password == expected_password:
            st.session_state.authenticated = True
            st.session_state.login_attempts = 0
            st.session_state.lockout_until = None
            st.rerun()
        else:
            st.session_state.login_attempts += 1
            if st.session_state.login_attempts >= 5:
                st.session_state.lockout_until = datetime.utcnow() + pd.Timedelta(seconds=60)
                st.session_state.login_attempts = 0
                st.error("Too many failed attempts. Locked for 60 seconds.")
            else:
                attempts_left = 5 - st.session_state.login_attempts
                st.error(f"Invalid credentials. {attempts_left} attempt(s) left before lockout.")

    st.stop()

def get_data_dir() -> Path:
    if getattr(sys, "frozen", False):
        base = Path(os.environ.get("LOCALAPPDATA", Path.home() / "AppData" / "Local"))
    else:
        base = Path(os.environ.get("LOCALAPPDATA", Path.home() / "AppData" / "Local"))
    data_dir = base / "ATA_Tool"
    data_dir.mkdir(parents=True, exist_ok=True)
    return data_dir
DATA_DIR = get_data_dir()
PARAMETERS_JSON = str(DATA_DIR / "parameters.json")
EXPORT_XLSX = str(DATA_DIR / "ATA_Audit_Log.xlsx")
SUMMARY_COLUMNS = [
    "Evaluation ID",
    "Evaluation Date",
    "Audit Date",
    "Reaudit",
    "QA Name",
    "Auditor",
    "Call ID",
    "Call Duration",
    "Call Disposition",
    "Overall Score %",
    "Passed Points",
    "Failed Points",
    "Total Points",
    "Last Updated",
]
DETAILS_COLUMNS = [
    "Evaluation ID",
    "Evaluation Date",
    "Audit Date",
    "Reaudit",
    "QA Name",
    "Auditor",
    "Call ID",
    "Overall Score %",
    "Group",
    "Parameter",
    "Points",
    "Description",
    "Result",
    "Comment",
]


def _norm_col_name(name: str) -> str:
    return str(name).strip().lower().replace("_", " ")


def _standardize_columns(df: pd.DataFrame, expected_columns: list[str]) -> pd.DataFrame:
    if df is None:
        df = pd.DataFrame()
    source = df.copy()
    rename_map = {_norm_col_name(col): col for col in source.columns}
    output = pd.DataFrame()
    for expected in expected_columns:
        actual = rename_map.get(_norm_col_name(expected))
        output[expected] = source[actual] if actual is not None else ""
    return output


def _rewrite_google_worksheet(ws, df: pd.DataFrame, expected_columns: list[str]) -> None:
    out = _standardize_columns(df, expected_columns)
    ws.clear()
    ws.append_row(expected_columns)
    if not out.empty:
        ws.append_rows(out.values.tolist())


def connect_google_sheet():
    scopes = ["https://www.googleapis.com/auth/spreadsheets"]
    creds_dict = dict(st.secrets["gcp_service_account"])
    # Convert escaped newlines into real newlines
    creds_dict["private_key"] = creds_dict["private_key"].replace("\\n", "\n")
    creds = Credentials.from_service_account_info(
        creds_dict,
        scopes=scopes
    )
    client = gspread.authorize(creds)
    sheet = client.open_by_key("15Rpo0vjtRVn9lUNXYcm4l2TRbmQOEbjB3qf9JQVpGPA")
    return sheet


def read_google_summary():
    sheet = connect_google_sheet()
    ws = sheet.worksheet("Summary")
    data = ws.get_all_records()
    df = _standardize_columns(pd.DataFrame(data), SUMMARY_COLUMNS)
    if "Evaluation ID" in df.columns:
        df["Evaluation ID"] = df["Evaluation ID"].astype(str).str.strip()
    for col in ["Passed Points", "Failed Points", "Total Points", "Overall Score %"]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)
    return df


def read_google_details():
    sheet = connect_google_sheet()
    ws = sheet.worksheet("Details")
    data = ws.get_all_records()
    df = _standardize_columns(pd.DataFrame(data), DETAILS_COLUMNS)
    if "Evaluation ID" in df.columns:
        df["Evaluation ID"] = df["Evaluation ID"].astype(str).str.strip()
    return df
DAMAC_TITLE = "DAMAC Properties"
DAMAC_SUB1 = "Quality Assurance"
DAMAC_SUB2 = "Telesales Division"
APP_NAME = "ATA Audit the Auditor"
LOGO_URL = "https://images.ctfassets.net/zoq5l15g49wj/2qCUAnuJ9WgJiGNxmTXkUa/0505928e3d7060e1bc0c3195d7def448/damac-gold.svg?fm=webp&w=200&h=202&fit=pad&q=60"
LOGIN_LOGO_URL = "https://www.moinsgroup.com/wp-content/uploads/2025/01/Damac_logo.svg-1-1024x256.png"
COOKIE_AUTH_KEY = "ata_auth"
COOKIE_REMEMBER_DAYS = 30


def cookie_expiry(days: int = COOKIE_REMEMBER_DAYS) -> datetime:
    return datetime.utcnow() + pd.Timedelta(days=days)


THEME_CONFIGS = {
    "light": {
        "main_bg": "#ffffff",
        "card_bg": "#F5F5DC",
        "title": "#CEAE72",
        "text": "#000000",
        "button_bg": "#F5F5DC",
        "button_text": "#CEAE72",
        "sidebar_text": "#000000",
        "group_bg": "#F5F5DC",
        "group_text": "#000000",
        "grid": "#0b1f3a",
        "border": "#CEAE72",
        "secondary": "#e7dcc2",
        "text_muted": "#374151",
        "card_shadow": "rgba(11,31,58,0.15)",
        "nav_bg": "#F5F5DC",
        "input_bg": "#f3f4f6",
        "pie_alt": ["#1e3a8a", "#CEAE72", "#2e5cb8", "#9c8a66", "#365ca1", "#c7a86e"],
    },
    "dark": {
        "main_bg": "#0b1f3a",
        "card_bg": "#08162a",
        "title": "#CEAE72",
        "text": "#ffffff",
        "button_bg": "#08162a",
        "button_text": "#CEAE72",
        "sidebar_text": "#ffffff",
        "group_bg": "#08162a",
        "group_text": "#CEAE72",
        "grid": "#b79f79",
        "border": "#CEAE72",
        "secondary": "#112843",
        "text_muted": "#cbd5e1",
        "card_shadow": "rgba(0,0,0,0.35)",
        "nav_bg": "#08162a",
        "input_bg": "#0d2748",
        "pie_alt": ["#CEAE72", "#4d6da8", "#8bb0f2", "#8aa16a", "#f2c97d", "#7094d1"],
    },
}


def get_active_theme() -> dict:
    preference = st.session_state.get("theme_mode", "system")
    if preference == "system":
        system_base = (st.get_option("theme.base") or "light").lower()
        mode = "dark" if system_base == "dark" else "light"
    else:
        mode = preference
    theme = THEME_CONFIGS[mode].copy()
    theme["mode"] = mode
    return theme


def apply_theme_css(theme: dict):
    is_dark = theme["mode"] == "dark"
    label_color = "#E6C77D" if is_dark else theme["text"]
    expander_color = "#F2D48A" if is_dark else theme["group_text"]
    table_header_bg = "#112843" if is_dark else theme["group_bg"]
    table_header_text = "#F5E6B3" if is_dark else theme["group_text"]
    table_body_text = "#FFFFFF" if is_dark else theme["text"]
    stat_hover_shadow = "rgba(0,0,0,0.35)" if is_dark else "rgba(0,0,0,0.12)"
    dashboard_hover_shadow = "rgba(0,0,0,0.35)" if is_dark else "rgba(0,0,0,0.18)"
    eval_btn_text = "#000000" if not is_dark else theme["button_text"]

    st.markdown(
        f"""
    <style>
    :root {{
        --primary: {theme['main_bg']};
        --secondary: {theme['secondary']};
        --accent-gold: {theme['title']};
        --bg-main: {theme['main_bg']};
        --bg-card: {theme['card_bg']};
        --text-main: {theme['text']};
        --text-muted: {theme['text_muted']};
        --border: {theme['border']};
        --btn-bg: {theme['button_bg']};
        --btn-text: {theme['button_text']};
        --group-bg: {theme['group_bg']};
        --group-text: {theme['group_text']};
        --grid: {theme['grid']};
        --sidebar-text: {theme['sidebar_text']};
        --card-shadow: {theme['card_shadow']};
        --nav-bg: {theme['nav_bg']};
        --input-bg: {theme['input_bg']};
    }}

    html, body, .stApp, [data-testid="stAppViewContainer"], [data-testid="stMainBlockContainer"] {{
        background: var(--bg-main) !important;
        color: var(--text-main) !important;
    }}

    .ata-card, .ata-title-card, .ata-nav-card, .stat-card {{
        background: var(--bg-card) !important;
        border: 1px solid var(--border) !important;
        border-radius: 16px;
        box-shadow: 0 10px 22px var(--card-shadow);
    }}

    .ata-title-card {{ padding: 1rem 1.2rem; margin: 0.2rem 0 1rem 0; }}
    .ata-title-card .title {{ color: var(--accent-gold); font-size: 2.15rem; font-weight: 800; margin: 0; }}
    .ata-title-card .subtitle {{ color: var(--text-main); margin-top: 0.35rem; font-size: 0.98rem; }}

    h1, h2, h3, h4, h5, h6, .ata-hero .t1, .view-detail-title, .view-score {{ color: var(--accent-gold) !important; }}

    .ata-hero {{
        padding: 30px;
        border-radius: 16px;
        background: var(--bg-card);
        border: 1px solid var(--border);
        box-shadow: 0 10px 22px var(--card-shadow);
        margin-bottom: 18px;
    }}
    .ata-hero .t1 {{ font-size: 34px; font-weight: 800; line-height: 1.2; color: var(--accent-gold); }}
    .ata-hero .t2 {{ font-size: 19px; font-weight: 500; line-height: 1.6; margin-top: 12px; color: var(--text-main); }}

    @media (max-width: 768px) {{
        .ata-hero .t1 {{ font-size: 27px; }}
        .ata-hero .t2 {{ font-size: 16px; }}
    }}

    .stat-card {{
        min-height: 148px;
        display: flex;
        flex-direction: column;
        justify-content: center;
        align-items: center;
        text-align: center;
        padding: 18px;
        transition: all 0.25s ease;
        cursor: pointer;
    }}
    .action-card [data-testid="stButton"] > button {{ transition: all 0.25s ease; cursor: pointer; }}
    .stat-card:hover, .action-card [data-testid="stButton"] > button:hover {{
        transform: translateY(-6px);
        box-shadow: 0 14px 30px {stat_hover_shadow};
    }}

    .stat-val {{ font-size: 2rem; font-weight: 800; color: var(--accent-gold); }}
    .stat-label {{ font-size: 1.05rem; color: var(--text-main); }}

    [data-testid="stSidebar"] {{ background: var(--bg-main) !important; border-right: 1px solid var(--border); }}
    [data-testid="stSidebar"] * {{ color: var(--sidebar-text) !important; }}

    .ata-nav-card {{ padding: 0.9rem 0.8rem; background: var(--nav-bg) !important; border: 1px solid var(--border) !important; margin-top: 0.8rem; margin-bottom: 0.9rem; }}
    .sidebar-nav-btn [data-testid="stButton"] > button {{ transition: all 0.25s ease; }}
    .sidebar-nav-btn [data-testid="stButton"] > button:hover {{ transform: translateY(-4px); box-shadow: 0 12px 24px {stat_hover_shadow}; }}

    .stButton>button,
    .stDownloadButton>button,
    .stForm [data-testid="stFormSubmitButton"]>button {{

        width: 100% !important;

        height: 60px !important;
        min-height: 60px !important;
        max-height: 60px !important;

        border-radius: 12px;
        border: 1px solid var(--border) !important;
        background: var(--btn-bg) !important;
        color: var(--btn-text) !important;

        font-weight: 700;
        font-size: 16px;

        display: flex;
        align-items: center;
        justify-content: center;

        margin: 0 !important;
    }}

    .eval-action-row [data-testid="stFormSubmitButton"] > button {{ color: {eval_btn_text} !important; }}

    .action-card {{
        width: 100%;
        height: 60px;
        display: flex;
        align-items: center;
    }}

    .view-action-grid {{
        display: grid;
        grid-template-columns: repeat(3, 1fr);
        gap: 12px;
        margin-top: 8px;
    }}
    .view-action-grid .view-action-cell {{
        min-width: 0;
    }}
    .view-action-grid .stButton>button,
    .view-action-grid .stDownloadButton>button,
    .view-action-grid .stForm [data-testid="stFormSubmitButton"]>button {{
        height: 46px !important;
        min-height: 46px !important;
        max-height: 46px !important;
        font-size: 14px !important;
        font-weight: 600 !important;
        padding: 0 12px !important;
    }}
    .view-action-grid iframe {{
        width: 100% !important;
        height: 46px !important;
        border: 0 !important;
    }}

    .stButton>button:hover, .stDownloadButton>button:hover, .stForm [data-testid="stFormSubmitButton"]>button:hover {{ transform: translateY(-4px); box-shadow: 0 12px 24px {stat_hover_shadow}; filter: brightness(1.06); }}

    .dashboard-action-btn [data-testid="stDownloadButton"] > button {{ transition: all 0.25s ease; cursor: pointer; }}
    .dashboard-action-btn [data-testid="stDownloadButton"] > button:hover {{ transform: translateY(-5px); box-shadow: 0 12px 28px {dashboard_hover_shadow}; }}
    .dashboard-action-btn [data-testid="stDownloadButton"] > button:active {{ transform: translateY(-2px); box-shadow: 0 6px 14px rgba(0,0,0,0.15); }}

    .styled-table {{ width: 100%; border-collapse: collapse; margin: 10px 0; font-size: 14px; }}
    .styled-table th {{ text-align:left; padding:12px 15px; background: var(--group-bg); color: var(--group-text); border:1px solid var(--grid); }}
    .styled-table td {{ padding:10px 15px; border:1px solid var(--grid); color: var(--text-main); }}
    .eval-details-table {{ width:100%; border-collapse:collapse; margin:10px 0; font-size:14px; }}
    .eval-details-table th, .eval-details-table td {{ padding:10px 12px; border:1px solid var(--border) !important; }}
    .summary-records-table {{ width:100%; border-collapse:collapse; border:1px solid var(--border); }}
    .summary-records-table th {{ background:#112843; color:#F2D48A; font-weight:600; border:1px solid var(--border); padding:10px 12px; }}
    .summary-records-table td {{ background:var(--bg-card); color:#FFFFFF; border:1px solid var(--border); padding:10px 12px; }}
    .summary-records-table tr:hover td {{ background:#132d4a; }}

    .dark-breakdown-table {{ width:100%; border-collapse:collapse; border:1px solid var(--border); }}
    .dark-breakdown-table th {{ background:#112843; color:#F5E6B3; border:1px solid var(--border); padding:10px 12px; text-align:left; }}
    .dark-breakdown-table td {{ background:var(--bg-card); color:#E5E7EB; border:1px solid var(--border); padding:9px 12px; }}

    [data-testid="stDataFrame"], [data-testid="stTable"] {{ border: 1px solid var(--border); border-radius: 12px; }}
    [data-testid="stDataFrame"] * {{ color: {table_body_text} !important; }}
    [data-testid="stDataEditor"] [role="columnheader"], [data-testid="stDataFrame"] thead th {{
        background: {table_header_bg} !important;
        color: {table_header_text} !important;
        border-color: var(--border) !important;
        text-align: left !important;
    }}
    [data-testid="stDataFrame"] thead th div {{
        justify-content: flex-start !important;
        text-align: left !important;
    }}
    [data-testid="stDataEditor"] [role="gridcell"] {{ color: {table_body_text} !important; }}

    label, .stSelectbox label, .stTextInput label, .stDateInput label {{ color: {label_color} !important; font-weight: 600 !important; }}
    [data-testid="stExpander"] summary p {{ color: {expander_color} !important; font-weight: 700 !important; }}

    .login-logo-card {{
        background: #0b1f3a;
        border: 2px solid #C9A227;
        border-radius: 18px;
        box-shadow: 0 8px 20px rgba(0,0,0,0.08);
        padding: 34px;
        text-align: center;
        margin-bottom: 14px;
    }}
    .login-form-shell {{ max-width: 480px; margin: 0 auto; }}
    .login-form-shell [data-testid="stForm"] {{ border: 0 !important; padding: 0 !important; background: transparent !important; }}

    .stTextInput input, .stDateInput input, .stSelectbox div[data-baseweb="select"] > div, .stTextArea textarea {{
        background: var(--input-bg) !important; color: var(--text-main) !important; border: 1px solid var(--border) !important;
    }}
    </style>
    """,
        unsafe_allow_html=True,
    )


def get_chart_theme() -> dict:
    theme = get_active_theme()
    return {
        "bg": theme["card_bg"],
        "grid": theme["grid"],
        "text": theme["text"],
        "title": theme["title"],
        "primary": "#1e3a8a" if theme["mode"] == "light" else "#CEAE72",
        "accent": "#CEAE72" if theme["mode"] == "light" else "#8bb0f2",
        "fail": "#ef4444",
        "pass": "#10b981",
        "pie_alt": theme["pie_alt"],
        "border": theme["border"],
    }


def style_chart(ax, theme: dict) -> None:
    ax.set_facecolor(theme["bg"])
    for spine in ax.spines.values():
        spine.set_color(theme["grid"])
    ax.tick_params(colors=theme["text"])
    ax.title.set_color(theme["title"])
    if ax.xaxis.label:
        ax.xaxis.label.set_color(theme["text"])
    if ax.yaxis.label:
        ax.yaxis.label.set_color(theme["text"])


def apply_base_css() -> None:
    st.markdown(
        """
<style>
.block-container { padding-top: 1rem; font-family: "Candara", "Segoe UI", sans-serif; }
.stApp, .stMarkdown, .stTextInput, .stSelectbox, .stDataEditor, .stButton, .stTable, .stDataFrame {
  font-family: "Candara", "Segoe UI", sans-serif;
}
.ata-hero.left-align { text-align:left; }

.login-logo {display:flex; justify-content:center; margin-bottom:12px;} 
.login-title {text-align:center; font-weight:800; margin-bottom:0;} 
.login-note {text-align:center; font-size:12px; margin-top:8px;} 
.login-extra {text-align:center; font-size:12px; margin-top:4px;} 
.page-spacer { height: 0.4rem; }
</style>
""",
        unsafe_allow_html=True,
    )


def render_title_card(title: str, subtitle: str = "") -> None:
    sub_html = f"<div class='subtitle'>{subtitle}</div>" if subtitle else ""
    st.markdown(
        f"""
        <div class="ata-title-card">
            <p class="title">{title}</p>
            {sub_html}
        </div>
        """,
        unsafe_allow_html=True,
    )

# -------------------- RUBRIC (SHEET-ALIGNED) --------------------
ACCURACY_HEADER = "Accuracy of Scoring"
ACCURACY_SUBPARAMS = [
    "Call Opening (Readiness / energy)",
    "Call Opening 2 (Confirming lead source / meeting focused)",
    "Effective Probing / Qualifying client",
    "Accurate / Complete info",
    "Objection / Call Handling",
    "Soft Skills (active listening - building rapport)",
    "Positivity / professionality / politeness",
    "Call closure / meeting summarized",
    "Accurate Disposition",
    "Comment / notes",
    "Accurate Data inputs / shows",
    "WhatsApp message sent",
    "Took lead ownership",
    "Follow-up made properly",
]
EVALUATION_QUALITY_PARAMS = [
    ("Adherence to QA Guidelines", "Followed QA process and aligned with calibration standards"),
    ("Evidence & Notes", "Left a clear, specific and improvement-focused comment"),
    ("Objectivity & Fairness", "Evaluation is unbiased and fact-based"),
    ("Critical Error Identification", "Correct identification of fatal errors"),
    (
        "Evaluation Variety & Sample Coverage",
        "Evaluations cover a balanced mix of call durations and call types",
    ),
    ("Feedback Actionability", "Conducted coaching session on the call topic (If required)"),
    ("Timeliness & Completeness", "On track with the evaluations target SLA"),
]
DEFAULT_PARAMETERS = {
    "form_name": APP_NAME,
    "parameters": [
        {
            "Parameter": ACCURACY_HEADER,
            "Description": "Header (not scored) | Category total = 14 points",
            "Points": 0,
            "Group": "HEADER",
        },
        *[
            {
                "Parameter": p,
                "Description": "Accuracy of Scoring – Sub Parameter",
                "Points": 1,
                "Group": "ACCURACY_SUB",
            }
            for p in ACCURACY_SUBPARAMS
        ],
        *[
            {"Parameter": p, "Description": d, "Points": 1, "Group": "EVAL_QUALITY"}
            for (p, d) in EVALUATION_QUALITY_PARAMS
        ],
    ],
}
def ensure_parameters_file() -> None:
    if not os.path.exists(PARAMETERS_JSON):
        with open(PARAMETERS_JSON, "w", encoding="utf-8") as f:
            json.dump(DEFAULT_PARAMETERS, f, indent=2, ensure_ascii=False)
def load_parameters_df() -> pd.DataFrame:
    ensure_parameters_file()
    with open(PARAMETERS_JSON, "r", encoding="utf-8") as f:
        cfg = json.load(f)
    df = pd.DataFrame(cfg.get("parameters", []))
    if df.empty:
        df = pd.DataFrame(DEFAULT_PARAMETERS["parameters"])
    df["Result"] = "Pass"
    df["Comment"] = ""
    for c in ["Parameter", "Points", "Description", "Result", "Comment", "Group"]:
        if c not in df.columns:
            df[c] = "" if c in ("Description", "Comment", "Group") else 1
    df["Points"] = pd.to_numeric(df["Points"], errors="coerce").fillna(0).astype(int)
    return df[["Group", "Parameter", "Points", "Description", "Result", "Comment"]].copy()
def normalize_details_df(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty or "Parameter" not in df.columns:
        return load_parameters_df()
    for col in ["Group", "Parameter", "Points", "Description", "Result", "Comment"]:
        if col not in df.columns:
            if col in ("Description", "Comment", "Group"):
                df[col] = ""
            elif col == "Result":
                df[col] = "Pass"
            else:
                df[col] = 1
    return df[["Group", "Parameter", "Points", "Description", "Result", "Comment"]].copy()
def format_date(value) -> str:
    if value in ("", None):
        return ""
    try:
        parsed = pd.to_datetime(value, dayfirst=True)
        if pd.isna(parsed):
            return ""
        return parsed.strftime("%d/%m/%Y")
    except Exception:
        return str(value)
def norm_id(x) -> str:
    if x is None:
        return ""
    try:
        if pd.isna(x):
            return ""
    except Exception:
        pass
    s = str(x).strip()
    if s.lower() in ("nan", "none"):
        return ""
    return s
def safe_read_excel(path: str, sheet: str) -> pd.DataFrame:
    if not os.path.exists(path):
        return pd.DataFrame()
    try:
        return pd.read_excel(path, sheet_name=sheet)
    except Exception:
        return pd.DataFrame()
def next_evaluation_id(evaluation_date_str: str) -> str:
    df = read_google_summary()
    yyyymmdd = evaluation_date_str.replace("-", "")
    prefix = f"ATA-{yyyymmdd}-"
    if df.empty or "Evaluation ID" not in df.columns:
        return f"{prefix}0001"
    existing = df["Evaluation ID"].astype(str)
    existing = existing[existing.str.startswith(prefix)]
    if existing.empty:
        return f"{prefix}0001"

    def _seq(x: str) -> int:
        try:
            return int(str(x).split("-")[-1])
        except Exception:
            return 0

    max_seq = max(existing.apply(_seq).tolist() + [0])
    return f"{prefix}{max_seq + 1:04d}"
def write_formatted_report(
    record: dict,
    filename: str,
    summary_df: pd.DataFrame | None = None,
    details_df: pd.DataFrame | None = None,
) -> None:
    if not os.path.exists(filename):
        return
    wb = load_workbook(filename)
    sheet_name = "Formatted Report"
    if sheet_name in wb.sheetnames:
        del wb[sheet_name]
    ws = wb.create_sheet(sheet_name)
    header_dark_blue = PatternFill(start_color="1F4E78", end_color="1F4E78", fill_type="solid")
    header_gray = PatternFill(start_color="C9C9C9", end_color="C9C9C9", fill_type="solid")
    header_light_gray = PatternFill(start_color="E5E5E5", end_color="E5E5E5", fill_type="solid")
    white_fill = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True)
    dark_font = Font(color="000000", bold=True)
    center = Alignment(horizontal="center", vertical="center", wrap_text=True)
    details_cols = [
        "Evaluation ID",
        "Evaluation Date",
        "Audit Date",
        "Reaudit",
        "QA Name",
        "Auditor",
        "Call ID",
        "Overall Score %",
    ]
    acc_cols = ACCURACY_SUBPARAMS
    eval_cols = [p for p, _ in EVALUATION_QUALITY_PARAMS]
    key_failed_col = "Key Failed Comments"
    all_cols = details_cols + acc_cols + eval_cols + [key_failed_col]
    details_start = 1
    details_end = details_start + len(details_cols) - 1
    acc_start = details_end + 1
    acc_end = acc_start + len(acc_cols) - 1
    eval_start = acc_end + 1
    eval_end = eval_start + len(eval_cols) - 1
    ws.row_dimensions[1].height = 10
    header_row = 2
    label_row = 3
    data_start_row = 4
    ws.merge_cells(start_row=header_row, start_column=details_start, end_row=header_row, end_column=details_end)
    ws.merge_cells(start_row=header_row, start_column=acc_start, end_row=header_row, end_column=acc_end)
    ws.merge_cells(start_row=header_row, start_column=eval_start, end_row=header_row, end_column=eval_end)
    ws.cell(row=header_row, column=details_start, value="Details").fill = header_dark_blue
    ws.cell(row=header_row, column=details_start).font = header_font
    ws.cell(row=header_row, column=details_start).alignment = center
    ws.cell(row=header_row, column=acc_start, value="ACCURACY_SUB").fill = header_gray
    ws.cell(row=header_row, column=acc_start).font = dark_font
    ws.cell(row=header_row, column=acc_start).alignment = center
    ws.cell(row=header_row, column=eval_start, value="EVAL_QUALITY").fill = header_light_gray
    ws.cell(row=header_row, column=eval_start).font = dark_font
    ws.cell(row=header_row, column=eval_start).alignment = center
    for idx, col_name in enumerate(all_cols, start=1):
        cell = ws.cell(row=label_row, column=idx, value=col_name)
        cell.alignment = center
        if idx <= details_end:
            cell.fill = header_dark_blue
            cell.font = header_font
        elif idx <= acc_end:
            cell.fill = header_gray
            cell.font = dark_font
        else:
            cell.fill = header_light_gray
            cell.font = dark_font
    summary = summary_df if summary_df is not None else safe_read_excel(filename, "Summary")
    details = details_df if details_df is not None else safe_read_excel(filename, "Details")
    if not summary.empty:
        summary = summary.dropna(axis=1, how="all")
    if summary.empty or details.empty:
        wb.save(filename)
        return
    if "Parameter" not in details.columns:
        wb.save(filename)
        return
    details_lookup = {
        eval_id: details[details["Evaluation ID"].astype(str).str.strip() == str(eval_id).strip()]
        for eval_id in summary["Evaluation ID"].dropna().unique().tolist()
    }
    for row_idx, row_data in enumerate(summary.to_dict(orient="records"), start=data_start_row):
        eval_id = row_data.get("Evaluation ID", "")
        details_values = [
            row_data.get("Evaluation ID", ""),
            format_date(row_data.get("Evaluation Date", "")),
            format_date(row_data.get("Audit Date", "")),
            row_data.get("Reaudit", ""),
            row_data.get("QA Name", ""),
            row_data.get("Auditor", ""),
            row_data.get("Call ID", ""),
            f"{row_data.get('Overall Score %', 0):.2f}%",
        ]
        detail_rows = details_lookup.get(eval_id, pd.DataFrame())
        if detail_rows.empty or "Parameter" not in detail_rows.columns:
            acc_values = ["" for _ in acc_cols]
            eval_values = ["" for _ in eval_cols]
        else:
            acc_values = []
            for param in acc_cols:
                match = detail_rows[detail_rows["Parameter"] == param]
                acc_values.append(match.iloc[0]["Result"] if not match.empty else "")
            eval_values = []
            for param in eval_cols:
                match = detail_rows[detail_rows["Parameter"] == param]
                eval_values.append(match.iloc[0]["Result"] if not match.empty else "")
        failed_comments = ""
        if not detail_rows.empty and "Result" in detail_rows.columns and "Comment" in detail_rows.columns:
            fail_rows = detail_rows[detail_rows["Result"].astype(str).str.strip().str.lower() == "fail"]
            comments = [str(c).strip() for c in fail_rows["Comment"].tolist() if str(c).strip() and str(c).strip().lower() != "nan"]
            failed_comments = " | ".join(dict.fromkeys(comments))
        values = details_values + acc_values + eval_values + [failed_comments]
        for idx, value in enumerate(values, start=1):
            cell = ws.cell(row=row_idx, column=idx, value=value)
            cell.fill = white_fill
            cell.alignment = center
        ws.row_dimensions[row_idx].height = 20
    ws.row_dimensions[header_row].height = 20
    ws.row_dimensions[label_row].height = 45
    for idx in range(1, len(all_cols) + 1):
        ws.column_dimensions[ws.cell(row=label_row, column=idx).column_letter].width = 18
    wb.save(filename)
def upsert_google_sheet(record: dict):
    rid = norm_id(record["evaluation_id"])
    if not rid:
        raise ValueError("Invalid Evaluation ID for upsert")

    sheet = connect_google_sheet()
    ws_summary = sheet.worksheet("Summary")
    ws_details = sheet.worksheet("Details")

    # Always read fresh data from Google Sheets to avoid stale-cache duplication.
    summary_existing = _standardize_columns(pd.DataFrame(ws_summary.get_all_records()), SUMMARY_COLUMNS)
    details_existing = _standardize_columns(pd.DataFrame(ws_details.get_all_records()), DETAILS_COLUMNS)

    if "Evaluation ID" in summary_existing.columns:
        summary_existing["Evaluation ID"] = summary_existing["Evaluation ID"].astype(str).str.strip()
        summary_existing = summary_existing[summary_existing["Evaluation ID"] != rid]
    if "Evaluation ID" in details_existing.columns:
        details_existing["Evaluation ID"] = details_existing["Evaluation ID"].astype(str).str.strip()
        details_existing = details_existing[details_existing["Evaluation ID"] != rid]

    summary_row = pd.DataFrame(
        [
            {
                "Evaluation ID": record["evaluation_id"],
                "Evaluation Date": record["evaluation_date"],
                "Audit Date": record["audit_date"],
                "Reaudit": record["reaudit"],
                "QA Name": record["qa_name"],
                "Auditor": record["auditor"],
                "Call ID": record["call_id"],
                "Call Duration": record["call_duration"],
                "Call Disposition": record["call_disposition"],
                "Overall Score %": record["overall_score"],
                "Passed Points": record["passed_points"],
                "Failed Points": record["failed_points"],
                "Total Points": record["total_points"],
                "Last Updated": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            }
        ]
    )
    details_df = normalize_details_df(record["details"])
    details_map = {
        "Evaluation ID": record["evaluation_id"],
        "Evaluation Date": record["evaluation_date"],
        "Audit Date": record["audit_date"],
        "Reaudit": record["reaudit"],
        "QA Name": record["qa_name"],
        "Auditor": record["auditor"],
        "Call ID": record["call_id"],
        "Overall Score %": record["overall_score"],
        "Points": details_df.get("Points", 1),
        "Description": details_df.get("Description", ""),
    }
    details_df = details_df.drop(columns=list(details_map.keys()), errors="ignore")
    for i, col in enumerate(details_map):
        details_df.insert(i, col, details_map[col])

    out_summary = pd.concat([summary_existing, summary_row], ignore_index=True)
    out_details = pd.concat([details_existing, details_df], ignore_index=True)

    # Idempotency/race protection: keep only the latest row(s) for this Evaluation ID.
    out_summary["Evaluation ID"] = out_summary["Evaluation ID"].astype(str).str.strip()
    out_summary = out_summary.drop_duplicates(subset=["Evaluation ID"], keep="last")
    out_details["Evaluation ID"] = out_details["Evaluation ID"].astype(str).str.strip()
    detail_dedupe_cols = [c for c in ["Evaluation ID", "Group", "Parameter"] if c in out_details.columns]
    if detail_dedupe_cols:
        out_details = out_details.drop_duplicates(subset=detail_dedupe_cols, keep="last")

    _rewrite_google_worksheet(ws_summary, out_summary, SUMMARY_COLUMNS)
    _rewrite_google_worksheet(ws_details, out_details, DETAILS_COLUMNS)

def delete_evaluation(eval_id: str) -> bool:
    rid = norm_id(eval_id)
    if not rid:
        return False
    summary_existing = read_google_summary()
    details_existing = read_google_details()
    if summary_existing.empty and details_existing.empty:
        return False

    summary_existing["_rid"] = summary_existing["Evaluation ID"].apply(norm_id)
    details_existing["_rid"] = details_existing["Evaluation ID"].apply(norm_id)

    before_summary = len(summary_existing)
    before_details = len(details_existing)

    summary_existing = summary_existing[summary_existing["_rid"] != rid].drop(columns=["_rid"], errors="ignore")
    details_existing = details_existing[details_existing["_rid"] != rid].drop(columns=["_rid"], errors="ignore")

    changed = (len(summary_existing) != before_summary) or (len(details_existing) != before_details)
    if not changed:
        return False

    sheet = connect_google_sheet()
    _rewrite_google_worksheet(sheet.worksheet("Summary"), summary_existing, SUMMARY_COLUMNS)
    _rewrite_google_worksheet(sheet.worksheet("Details"), details_existing, DETAILS_COLUMNS)
    st.cache_data.clear()
    return True

def compute_weighted_score(df: pd.DataFrame) -> dict:
    df = df.copy()
    df["Result"] = df["Result"].fillna("Pass")
    accuracy_rows = df[df["Group"] == "ACCURACY_SUB"]
    eval_quality_rows = df[df["Group"] == "EVAL_QUALITY"]
    accuracy_failed = (accuracy_rows["Result"] == "Fail").any()
    accuracy_points = 1
    accuracy_passed = 0 if accuracy_failed else 1
    eval_quality_total = int(len(eval_quality_rows))
    eval_quality_passed = int((eval_quality_rows["Result"] == "Pass").sum())
    total_points = accuracy_points + eval_quality_total
    passed_points = accuracy_passed + eval_quality_passed
    failed_points = total_points - passed_points
    score = round((passed_points / total_points) * 100, 2) if total_points else 0.0
    return {
        "score": score,
        "passed_points": passed_points,
        "failed_points": failed_points,
        "total_points": total_points,
    }
# -------------------- EXPORT HELPERS --------------------
def copy_html_to_clipboard_button(label: str, html_to_copy: str, key: str, theme: dict) -> None:
    html_b64 = base64.b64encode(html_to_copy.encode("utf-8")).decode("ascii")

    js = f"""
    <style>
      html, body {{ margin:0; padding:0; background:transparent; }}
      #btn-{key} {{
        width:100%;
        height:46px;
        border-radius:10px;
        border:1px solid {theme['border']};
        background:{theme['button_bg']};
        color:{theme['button_text']};
        font-weight:700;
        padding:0.5rem 1rem;
        cursor:pointer;
        transition: all 0.25s ease;
      }}
      #btn-{key}:hover {{
        transform: translateY(-4px);
        box-shadow: 0 12px 24px rgba(0,0,0,0.25);
        filter: brightness(1.05);
      }}
      #status-{key} {{
        margin-top: 4px;
        min-height: 16px;
        font-size: 12px;
        line-height: 1.2;
        color: {theme.get('button_text', '#000000')};
        text-align: center;
        opacity: 0;
        transition: opacity 0.2s ease;
      }}
      #status-{key}.show {{ opacity: 1; }}
      #status-{key}.err {{ color:{theme.get('fail', '#ef4444')}; }}
    </style>

    <button id="btn-{key}" type="button">{label}</button>
    <div id="status-{key}"></div>

    <script>
      const b64 = "{html_b64}";
      const statusEl = document.getElementById("status-{key}");

      function b64ToUtf8(b64Str) {{
        const bin = atob(b64Str);
        const bytes = new Uint8Array([...bin].map(ch => ch.charCodeAt(0)));
        return new TextDecoder("utf-8").decode(bytes);
      }}

      function copyRichHTML() {{
        const html = b64ToUtf8(b64);

        const container = document.createElement("div");
        container.innerHTML = html;
        container.contentEditable = true;
        container.style.position = "fixed";
        container.style.left = "-9999px";
        container.style.top = "0";
        document.body.appendChild(container);

        const range = document.createRange();
        range.selectNodeContents(container);
        const selection = window.getSelection();
        selection.removeAllRanges();
        selection.addRange(range);

        try {{
          const successful = document.execCommand("copy");
          if (successful) {{
            statusEl.innerText = "Copied successfully";
            statusEl.className = "show";
            setTimeout(() => {{ statusEl.innerText = ""; statusEl.className = ""; }}, 2000);
          }} else {{
            throw new Error("Copy failed");
          }}
        }} catch (err) {{
          statusEl.innerText = "Copy blocked by browser";
          statusEl.className = "err show";
          setTimeout(() => {{ statusEl.innerText = ""; statusEl.className = ""; }}, 2000);
        }}

        selection.removeAllRanges();
        document.body.removeChild(container);
      }}

      document.getElementById("btn-{key}").addEventListener("click", copyRichHTML);
    </script>
    """

    components.html(js, height=60)

def email_subject_text(record: dict) -> str:
    return f"ATA Evaluation | {record['evaluation_id']} | {record['qa_name']} | {format_date(record['audit_date'])}"

def email_html_inline(record: dict) -> str:
    def make_table(df):
        rows = []
        for _, r in df.iterrows():
            status = r["Result"]
            badge = (
                "background:#e9f7ef;color:#1f8f4a"
                if status == "Pass"
                else "#fde8e6;color:#d93025"
            )
            badge = f"{badge};padding:4px 8px;border-radius:12px;font-weight:bold;"
            comm = str(r["Comment"]) if str(r["Comment"]).lower() != "nan" else ""
            rows.append(
                "<tr>"
                f"<td style='padding:10px 12px;border-bottom:1px solid #e5e7eb;vertical-align:top;width:45%;'>{r['Parameter']}</td>"
                "<td style='padding:10px 12px;border-bottom:1px solid #e5e7eb;text-align:center;vertical-align:middle;width:15%;'>"
                f"<span style='{badge}'>{status}</span>"
                "</td>"
                f"<td style='padding:10px 12px;border-bottom:1px solid #e5e7eb;vertical-align:top;width:40%;'>{comm}</td>"
                "</tr>"
            )
        return (
            "<table role='presentation' cellpadding='0' cellspacing='0' style='width:100%;border-collapse:collapse;table-layout:fixed;font-size:13px;'>"
            "<tr style='background:#f8fafc;'>"
            "<th style='text-align:left;padding:10px 12px;border-bottom:1px solid #e5e7eb;width:45%;'>Parameter</th>"
            "<th style='text-align:center;padding:10px 12px;border-bottom:1px solid #e5e7eb;width:15%;'>Result</th>"
            "<th style='text-align:left;padding:10px 12px;border-bottom:1px solid #e5e7eb;width:40%;'>Comment</th>"
            "</tr>"
            f"{''.join(rows)}"
            "</table>"
        )
    det = record["details"]
    return f"""
    <table role="presentation" cellpadding="0" cellspacing="0" width="100%" style="margin:0;padding:0;border-collapse:collapse;background:#f3f4f6;">
      <tr>
        <td align="center" style="padding:24px 12px;">
          <table role="presentation" cellpadding="0" cellspacing="0" width="700" style="width:700px;max-width:700px;border-collapse:separate;border-spacing:0;background:#ffffff;border:1px solid #e5e7eb;border-radius:12px;box-shadow:0 6px 18px rgba(15,23,42,0.08);overflow:hidden;font-family:Candara, Segoe UI, sans-serif;color:#111827;">
            <tr>
              <td style="background:#0b1f3a;color:#ffffff;padding:12px 18px;">
                <div style="font-size:20px;line-height:1.25;font-weight:700;">{DAMAC_TITLE} | ATA Evaluation</div>
                <div style="font-size:12px;line-height:1.5;opacity:0.9;margin-top:2px;">{DAMAC_SUB1} | {DAMAC_SUB2}</div>
              </td>
            </tr>
            <tr>
              <td style="padding:16px 18px 6px 18px;">
                <table role="presentation" cellpadding="0" cellspacing="0" width="100%" style="width:100%;border-collapse:collapse;table-layout:fixed;font-size:13px;">
                  <tr>
                    <td style="padding:6px 8px;vertical-align:top;width:50%;"><span style="font-weight:700;color:#0b1f3a;">Evaluation ID:</span> {record['evaluation_id']}</td>
                    <td style="padding:6px 8px;vertical-align:top;width:50%;"><span style="font-weight:700;color:#0b1f3a;">Evaluation Date:</span> {record['evaluation_date']}</td>
                  </tr>
                  <tr>
                    <td style="padding:6px 8px;vertical-align:top;"><span style="font-weight:700;color:#0b1f3a;">QA Name:</span> {record['qa_name']}</td>
                    <td style="padding:6px 8px;vertical-align:top;"><span style="font-weight:700;color:#0b1f3a;">Auditor Name:</span> {record['auditor']}</td>
                  </tr>
                  <tr>
                    <td style="padding:6px 8px;vertical-align:top;"><span style="font-weight:700;color:#0b1f3a;">Audit Date:</span> {format_date(record['audit_date'])}</td>
                    <td style="padding:6px 8px;vertical-align:top;"><span style="font-weight:700;color:#0b1f3a;">Call ID:</span> {record['call_id']}</td>
                  </tr>
                  <tr>
                    <td style="padding:6px 8px;vertical-align:top;"><span style="font-weight:700;color:#0b1f3a;">Call Duration:</span> {record['call_duration']}</td>
                    <td style="padding:6px 8px;vertical-align:top;"><span style="font-weight:700;color:#0b1f3a;">Call Disposition:</span> {record['call_disposition']}</td>
                  </tr>
                </table>
              </td>
            </tr>
            <tr>
              <td style="padding:8px 18px 0 18px;">
                <table role="presentation" cellpadding="0" cellspacing="0" width="100%" style="border-collapse:collapse;">
                  <tr>
                    <td align="center" style="background:#f1f5f9;border:1px solid #e5e7eb;border-radius:8px;padding:10px 12px;">
                      <div style="font-size:12px;color:#475569;line-height:1.4;">Overall Score</div>
                      <div style="font-size:24px;font-weight:700;line-height:1.2;color:#0b1f3a;">{record['overall_score']:.2f}%</div>
                    </td>
                  </tr>
                </table>
              </td>
            </tr>
            <tr>
              <td style="padding:18px 18px 0 18px;">
                <div style="font-size:18px;font-weight:700;line-height:1.3;color:#0b1f3a;border-left:4px solid #0b1f3a;padding-left:10px;">Accuracy of Scoring</div>
              </td>
            </tr>
            <tr>
              <td style="padding:10px 18px 0 18px;">{make_table(det[det["Group"] == "ACCURACY_SUB"])}</td>
            </tr>
            <tr>
              <td style="padding:22px 18px 0 18px;">
                <div style="font-size:18px;font-weight:700;line-height:1.3;color:#0b1f3a;border-left:4px solid #0b1f3a;padding-left:10px;">Evaluation Quality</div>
              </td>
            </tr>
            <tr>
              <td style="padding:10px 18px 0 18px;">{make_table(det[det["Group"] == "EVAL_QUALITY"])}</td>
            </tr>
            <tr>
              <td style="padding:18px;">
                <table role="presentation" cellpadding="0" cellspacing="0" width="100%" style="border-collapse:collapse;">
                  <tr>
                    <td style="background:#f8fafc;border:1px solid #e5e7eb;border-radius:8px;padding:10px 12px;font-size:13px;line-height:1.4;"><span style="font-weight:700;color:#0b1f3a;">Reaudit Status:</span> <span style="font-weight:700;color:{'#d93025' if str(record['reaudit']).strip().lower() == 'yes' else '#1f8f4a'};">{record['reaudit']}</span></td>
                  </tr>
                </table>
              </td>
            </tr>
          </table>
        </td>
      </tr>
    </table>
    """
class PDFReport(FPDF):
    def header(self):
        self.set_font("Arial", "B", 10)
        self.set_text_color(11, 31, 58)
        self.cell(0, 8, f"{DAMAC_TITLE} | {DAMAC_SUB1} | {DAMAC_SUB2}", ln=True, align="C")
        self.ln(2)
    def footer(self):
        self.set_y(-15)
        self.set_font("Arial", "I", 8)
        self.set_text_color(128, 128, 128)
        self.cell(0, 10, f"Page {self.page_no()}", align="C")
def pdf_evaluation(record: dict) -> bytes:
    pdf = PDFReport()
    pdf.add_page()
    pdf.set_auto_page_break(auto=False, margin=12)
    pdf.set_font("Arial", "B", 14)
    pdf.cell(0, 12, "ATA Evaluation Report", ln=True, align="L")
    pdf.ln(2)
    pdf.set_font("Arial", "B", 9)
    pdf.set_fill_color(241, 245, 249)
    data = [
        ["Evaluation ID", record["evaluation_id"], "Evaluation Date", format_date(record["evaluation_date"])],
        ["QA Name", record["qa_name"], "Auditor Name", record["auditor"]],
        ["Audit Date", format_date(record["audit_date"]), "Call ID", record["call_id"]],
        ["Call Duration", record["call_duration"], "Call Disposition", record["call_disposition"]],
    ]
    def line_count(text: str, width: float) -> int:
        words = str(text).split()
        if not words:
            return 1
        lines = 1
        line = ""
        for word in words:
            test = f"{line} {word}".strip()
            if pdf.get_string_width(test) <= width:
                line = test
            else:
                lines += 1
                line = word
        return lines
    col_widths = [42, 42, 42, 64]
    def truncate_text(text: str, width: float) -> str:
        text = str(text)
        if pdf.get_string_width(text) <= width:
            return text
        while text and pdf.get_string_width(f"{text}...") > width:
            text = text[:-1]
        return f"{text}..." if text else ""
    row_height = 10
    for row in data:
        x = pdf.get_x()
        y = pdf.get_y()
        pdf.set_font("Arial", "B", 9)
        pdf.cell(col_widths[0], row_height, truncate_text(row[0], col_widths[0]), border=1, fill=True)
        pdf.set_font("Arial", "", 9)
        pdf.cell(col_widths[1], row_height, truncate_text(row[1], col_widths[1]), border=1)
        pdf.set_font("Arial", "B", 9)
        pdf.cell(col_widths[2], row_height, truncate_text(row[2], col_widths[2]), border=1, fill=True)
        pdf.set_font("Arial", "", 9)
        pdf.cell(col_widths[3], row_height, truncate_text(row[3], col_widths[3]), border=1)
        pdf.ln(row_height)
    pdf.ln(4)
    pdf.set_font("Arial", "B", 11)
    pdf.cell(
        0,
        10,
        f"Overall Score: {record['overall_score']:.2f}%",
        ln=True,
        align="C",
        border=1,
        fill=True,
    )
    pdf.ln(5)
    def draw_section(title, df):
        pdf.set_font("Arial", "B", 10)
        pdf.set_text_color(11, 31, 58)
        pdf.cell(0, 8, title, ln=True)
        pdf.set_text_color(0, 0, 0)
        pdf.set_font("Arial", "B", 8)
        pdf.set_fill_color(11, 31, 58)
        pdf.set_text_color(255, 255, 255)
        param_w, result_w, comment_w = 90, 20, 80
        pdf.cell(param_w, 8, " Parameter", border=1, fill=True)
        pdf.cell(result_w, 8, " Result", border=1, fill=True, align="C")
        pdf.cell(comment_w, 8, " Comment", border=1, fill=True)
        pdf.ln()
        pdf.set_font("Arial", "", 7)
        pdf.set_text_color(0, 0, 0)
        def truncate_section(text: str, width: float) -> str:
            text = str(text)
            if pdf.get_string_width(text) <= width:
                return text
            while text and pdf.get_string_width(f"{text}...") > width:
                text = text[:-1]
            return f"{text}..." if text else ""
        def wrap_lines(text: str, width: float) -> list[str]:
            words = str(text).split()
            if not words:
                return [""]
            lines = []
            line = ""
            for word in words:
                test = f"{line} {word}".strip()
                if pdf.get_string_width(test) <= width:
                    line = test
                else:
                    lines.append(line)
                    line = word
            if line:
                lines.append(line)
            return lines
        def ensure_space(height: float) -> None:
            if pdf.get_y() + height > (pdf.h - 12):
                pdf.add_page()
                pdf.set_font("Arial", "B", 10)
                pdf.set_text_color(11, 31, 58)
                pdf.cell(0, 8, title, ln=True)
                pdf.set_text_color(0, 0, 0)
                pdf.set_font("Arial", "B", 8)
                pdf.set_fill_color(11, 31, 58)
                pdf.set_text_color(255, 255, 255)
                pdf.cell(param_w, 8, " Parameter", border=1, fill=True)
                pdf.cell(result_w, 8, " Result", border=1, fill=True, align="C")
                pdf.cell(comment_w, 8, " Comment", border=1, fill=True)
                pdf.ln()
                pdf.set_font("Arial", "", 7)
                pdf.set_text_color(0, 0, 0)
        for _, r in df.iterrows():
            param, res, comm = str(r["Parameter"]), str(r["Result"]), str(r["Comment"])
            if comm.lower() == "nan":
                comm = ""
            comment_lines = wrap_lines(comm, comment_w - 2) if comm.strip() else [""]
            row_height = max(6, 4 * len(comment_lines))
            line_height = 4
            ensure_space(row_height + 2)
            start_x, start_y = pdf.get_x(), pdf.get_y()
            pdf.rect(start_x, start_y, param_w, row_height)
            pdf.rect(start_x + param_w, start_y, result_w, row_height)
            pdf.rect(start_x + param_w + result_w, start_y, comment_w, row_height)
            y_offset = (row_height - line_height) / 2
            pdf.set_xy(start_x + 1, start_y + y_offset)
            pdf.cell(param_w - 2, line_height, truncate_section(param, param_w - 2), border=0)
            pdf.set_xy(start_x + param_w, start_y + y_offset)
            if res == "Pass":
                pdf.set_text_color(31, 143, 74)
            elif res == "Fail":
                pdf.set_text_color(217, 48, 37)
            else:
                pdf.set_text_color(0, 0, 0)
            pdf.cell(result_w, line_height, truncate_section(res, result_w - 2), border=0, align="C")
            pdf.set_text_color(0, 0, 0)
            pdf.set_xy(start_x + param_w + result_w + 1, start_y + 1)
            pdf.multi_cell(comment_w - 2, line_height, "\n".join(comment_lines), border=0)
            pdf.set_xy(start_x, start_y + row_height)
        pdf.ln(5)
    det = record["details"]
    draw_section("Accuracy of Scoring", det[det["Group"] == "ACCURACY_SUB"])
    draw_section("Evaluation Quality", det[det["Group"] == "EVAL_QUALITY"])
    pdf.set_font("Arial", "B", 10)
    pdf.set_text_color(0, 0, 0)
    pdf.cell(34, 8, "Reaudit Status:", ln=0)
    if str(record['reaudit']).strip().lower() == "yes":
        pdf.set_text_color(217, 48, 37)
    else:
        pdf.set_text_color(31, 143, 74)
    pdf.cell(0, 8, f" {record['reaudit']}", ln=True)
    pdf.set_text_color(0, 0, 0)
    out = pdf.output(dest="S")
    return bytes(out) if isinstance(out, (bytes, bytearray)) else out.encode("latin-1")
# ================= NEW INTELLIGENCE LAYER =================
def compute_auditor_intelligence(summary_df: pd.DataFrame, details_df: pd.DataFrame) -> pd.DataFrame:
    cols = [
        "Auditor",
        "Avg Score",
        "Failure Rate",
        "Reaudit Ratio",
        "Volatility",
        "Repeat Failure Count",
    ]
    if summary_df is None or summary_df.empty:
        return pd.DataFrame(columns=cols)

    summary = summary_df.copy()
    if "Auditor" not in summary.columns:
        return pd.DataFrame(columns=cols)

    summary["Auditor"] = summary["Auditor"].fillna("Unknown").astype(str).str.strip()
    summary["Evaluation Date"] = pd.to_datetime(summary.get("Evaluation Date"), errors="coerce")
    for col in ["Overall Score %", "Failed Points", "Total Points"]:
        if col in summary.columns:
            summary[col] = pd.to_numeric(summary[col], errors="coerce").fillna(0)
        else:
            summary[col] = 0

    summary["_reaudit_yes"] = summary.get("Reaudit", "").astype(str).str.strip().str.lower().eq("yes")
    grp = summary.groupby("Auditor", dropna=False)
    agg = grp.agg(
        avg_score=("Overall Score %", "mean"),
        fail_points=("Failed Points", "sum"),
        total_points=("Total Points", "sum"),
        reaudit_ratio=("_reaudit_yes", "mean"),
        volatility=("Overall Score %", "std"),
    )
    agg["Failure Rate"] = (agg["fail_points"] / agg["total_points"].replace(0, pd.NA)).fillna(0) * 100
    agg["Reaudit Ratio"] = agg["reaudit_ratio"].fillna(0) * 100
    agg["Volatility"] = agg["volatility"].fillna(0)

    repeat_counts = pd.Series(dtype=float)
    if details_df is not None and not details_df.empty:
        det = details_df.copy()
        det["Auditor"] = det.get("Auditor", "Unknown").fillna("Unknown").astype(str).str.strip()
        det["Parameter"] = det.get("Parameter", "").fillna("").astype(str).str.strip()
        det["Result"] = det.get("Result", "").fillna("").astype(str).str.strip()
        det["Evaluation Date"] = pd.to_datetime(det.get("Evaluation Date"), errors="coerce")
        latest_dt = det["Evaluation Date"].max()
        if pd.notna(latest_dt):
            window_start = latest_dt - pd.Timedelta(days=30)
            det = det[det["Evaluation Date"] >= window_start]
        fail_det = det[(det["Result"].str.lower() == "fail") & (det["Parameter"] != "")]
        if not fail_det.empty:
            by_param = fail_det.groupby(["Auditor", "Parameter"], dropna=False).size().reset_index(name="fail_count")
            repeated = by_param[by_param["fail_count"] >= 3]
            if not repeated.empty:
                repeat_counts = repeated.groupby("Auditor").size()

    out = pd.DataFrame(index=agg.index)
    out["Avg Score"] = agg["avg_score"].fillna(0)
    out["Failure Rate"] = agg["Failure Rate"].fillna(0)
    out["Reaudit Ratio"] = agg["Reaudit Ratio"].fillna(0)
    out["Volatility"] = agg["Volatility"].fillna(0)
    out["Repeat Failure Count"] = repeat_counts.reindex(out.index).fillna(0).astype(int)
    out = out.reset_index().rename(columns={"index": "Auditor"})
    return out[cols]


def compute_risk_flags(auditor_df: pd.DataFrame, details_df: pd.DataFrame) -> pd.DataFrame:
    cols = ["Auditor", "Risk Points", "Risk Level", "QA Intervention Required", "Coaching Required"]
    if auditor_df is None or auditor_df.empty:
        return pd.DataFrame(columns=cols)

    aud = auditor_df.copy()
    aud["Auditor"] = aud["Auditor"].fillna("Unknown").astype(str).str.strip()
    for c in ["Avg Score", "Repeat Failure Count", "Reaudit Ratio"]:
        aud[c] = pd.to_numeric(aud.get(c, 0), errors="coerce").fillna(0)

    critical_fail = pd.Series(0, index=aud["Auditor"].values, dtype=int)
    if details_df is not None and not details_df.empty:
        det = details_df.copy()
        det["Auditor"] = det.get("Auditor", "Unknown").fillna("Unknown").astype(str).str.strip()
        det["Parameter"] = det.get("Parameter", "").fillna("").astype(str).str.strip().str.lower()
        det["Result"] = det.get("Result", "").fillna("").astype(str).str.strip().str.lower()
        crit = det[(det["Parameter"] == "critical error identification") & (det["Result"] == "fail")]
        if not crit.empty:
            critical_fail = crit.groupby("Auditor").size().reindex(aud["Auditor"]).fillna(0).astype(int).values

    points = (
        (aud["Avg Score"] < 85).astype(int) * 2
        + (aud["Repeat Failure Count"] >= 1).astype(int) * 2
        + (pd.Series(critical_fail, index=aud.index) > 0).astype(int) * 3
        + (aud["Reaudit Ratio"] > 30).astype(int) * 1
    )

    out = pd.DataFrame({"Auditor": aud["Auditor"], "Risk Points": points.astype(int)})
    out["Risk Level"] = out["Risk Points"].map(lambda x: "Low" if x <= 1 else ("Moderate" if x <= 3 else "High"))
    out["QA Intervention Required"] = out["Risk Level"].eq("High")
    out["Coaching Required"] = out["Risk Level"].isin(["Moderate", "High"])
    return out[cols]


def compute_health_index(auditor_df: pd.DataFrame, details_df: pd.DataFrame) -> pd.DataFrame:
    cols = ["Auditor", "Health Index", "Health Classification", "Critical Fail Rate"]
    if auditor_df is None or auditor_df.empty:
        return pd.DataFrame()

    if details_df is None:
        details_df = pd.DataFrame()

    aud = auditor_df.copy()
    aud["Auditor"] = aud["Auditor"].fillna("Unknown").astype(str).str.strip()
    for c in ["Avg Score", "Failure Rate", "Reaudit Ratio"]:
        aud[c] = pd.to_numeric(aud.get(c, 0), errors="coerce").fillna(0)

    critical_fail_rate = pd.Series(0.0, index=aud["Auditor"])
    if not details_df.empty:
        det = details_df.copy()
        det["Auditor"] = det.get("Auditor", "Unknown").fillna("Unknown").astype(str).str.strip()
        det["Parameter"] = det.get("Parameter", "").fillna("").astype(str).str.strip().str.lower()
        det["Result"] = det.get("Result", "").fillna("").astype(str).str.strip().str.lower()
        crit_all = det[det["Parameter"] == "critical error identification"]
        if not crit_all.empty:
            crit_rate = (crit_all["Result"].eq("fail").groupby(crit_all["Auditor"]).mean() * 100)
            critical_fail_rate = crit_rate.reindex(aud["Auditor"]).fillna(0)

    out = pd.DataFrame({"Auditor": aud["Auditor"]})
    out["Critical Fail Rate"] = pd.to_numeric(
        critical_fail_rate, errors="coerce"
    ).fillna(0)
    out["Health Index"] = (
        0.40 * aud["Avg Score"]
        + 0.25 * (100 - aud["Failure Rate"])
        + 0.20 * (100 - out["Critical Fail Rate"])
        + 0.15 * (100 - aud["Reaudit Ratio"])
    ).clip(0, 100)
    out["Health Classification"] = out["Health Index"].map(
        lambda x: "Excellent" if x >= 90 else ("Stable" if x >= 75 else ("Watchlist" if x >= 60 else "High Risk"))
    )
    return out[cols]


def generate_coaching_summary(evaluation_record: dict, auditor_metrics: dict | pd.Series | None) -> str:
    if evaluation_record is None:
        return "No evaluation data available."

    if auditor_metrics is None:
        auditor_metrics = {}

    details = evaluation_record.get("details", pd.DataFrame())
    if details is None or not isinstance(details, pd.DataFrame) or details.empty:
        return "No parameter details available to generate coaching summary."

    det = details.copy()
    det["Parameter"] = det.get("Parameter", "").fillna("").astype(str)
    det["Result"] = det.get("Result", "").fillna("").astype(str)
    det["Comment"] = det.get("Comment", "").fillna("").astype(str)
    failed_df = det.loc[det["Result"].str.lower() == "fail", ["Parameter", "Comment"]].copy()

    metrics = dict(auditor_metrics) if isinstance(auditor_metrics, (dict, pd.Series)) else {}
    risk_level = str(metrics.get("Risk Level", "Low"))
    follow_up = "7 days" if risk_level == "High" else ("14 days" if risk_level == "Moderate" else "Monitor next cycle")

    failed_params = [str(p).strip() for p in failed_df["Parameter"].tolist() if str(p).strip()]
    governance_gaps = "\n".join([f"- {p}" for p in failed_params]) if failed_params else "- No governance gaps identified in this cycle."

    if risk_level == "High":
        risk_observation = "Elevated governance exposure detected. Immediate intervention required."
    elif risk_level == "Moderate":
        risk_observation = "Moderate governance variance observed. Focused coaching required."
    else:
        risk_observation = "Performance within governance tolerance. Continue monitoring."

    actions = []
    for _, row in failed_df.drop_duplicates(subset=["Parameter", "Comment"]).iterrows():
        parameter = str(row.get("Parameter", "")).strip()
        comment = str(row.get("Comment", "")).strip()
        if not parameter:
            continue
        if comment:
            actions.append(
                f"- For '{parameter}', evaluation indicates: '{comment}'. Action: Conduct focused recalibration, reinforce governance control, and validate consistency in next two evaluations."
            )
        else:
            actions.append(
                f"- For '{parameter}', conduct targeted governance coaching and revalidate scoring discipline."
            )
        if parameter.lower() == "critical error identification":
            actions.append(
                "- Immediate governance escalation required. Align with QA leadership for recalibration review."
            )
    recommended_action_plan = "\n".join(dict.fromkeys(actions)) if actions else "- Continue governance monitoring and sustain current control discipline."

    return (
        f"Senior QA Governance Coaching\n"
        f"Evaluation ID: {evaluation_record.get('evaluation_id', '')}\n"
        f"Auditor Under Review: {evaluation_record.get('auditor', '')}\n"
        f"Risk Level: {risk_level}\n\n"
        f"Governance Gaps\n{governance_gaps}\n\n"
        f"Risk Observation\n{risk_observation}\n\n"
        f"Recommended Action Plan\n{recommended_action_plan}\n\n"
        f"Follow-Up Timeline\n- {follow_up}"
    )

# -------------------- DASHBOARD LOGIC --------------------
def build_dashboard_figs(summary: pd.DataFrame | None = None, details: pd.DataFrame | None = None):
    if summary is None or details is None:
        summary = read_google_summary()
        details = read_google_details()
    if summary.empty or details.empty:
        return (None,) * 10 + (summary, details)
    for col in ["Passed Points", "Failed Points", "Total Points", "Overall Score %"]:
        if col in summary.columns:
            summary[col] = pd.to_numeric(summary[col], errors="coerce").fillna(0)
    summary["Evaluation Date"] = pd.to_datetime(summary.get("Evaluation Date"), errors="coerce")
    summary["Month"] = summary["Evaluation Date"].dt.to_period("M")
    summary["Failure Rate"] = summary.apply(
        lambda r: (r["Failed Points"] / r["Total Points"]) if r["Total Points"] else 0,
        axis=1,
    )
    theme = get_chart_theme()

    def add_bar_labels(ax):
        heights = [patch.get_height() for patch in ax.patches]
        for patch in ax.patches:
            value = patch.get_height()
            ax.annotate(
                f"{value:.1f}" if isinstance(value, float) else f"{value}",
                (patch.get_x() + patch.get_width() / 2, value),
                ha="center",
                va="bottom",
                xytext=(0, 6),
                textcoords="offset points",
                fontsize=9,
                color=theme["text"],
            )
    # 1. Trend Chart (Failure Rate)
    trend = summary.groupby("Month")["Failure Rate"].mean().sort_index()
    trend_x = trend.index.to_timestamp()
    fig_trend, ax = plt.subplots(figsize=(6, 4))
    ax.plot(trend_x, trend.values * 100, marker="o", color=theme["primary"], linewidth=2, markersize=6)
    ax.fill_between(trend_x, trend.values * 100, color=theme["primary"], alpha=0.15)
    ax.set_title("Failure Rate Trend (%)", fontweight="bold", fontsize=11)
    ax.grid(True, alpha=0.25, color=theme["grid"])
    ax.set_xticks(trend_x)
    ax.set_xticklabels([d.strftime("%b-%y") for d in trend_x])
    for x, y in zip(trend_x, trend.values * 100):
        ax.annotate(f"{y:.1f}%", (x, y), textcoords="offset points", xytext=(0, 10), ha="center", color=theme["accent"], fontweight="bold")
    style_chart(ax, theme)
    fig_trend.patch.set_facecolor(theme["bg"])
    plt.tight_layout()
    # 2. Heatmap
    fail_rows = details[details.get("Result", "") == "Fail"].copy()
    heat = pd.DataFrame()
    if not fail_rows.empty:
        fail_rows["Date Label"] = pd.to_datetime(
            fail_rows.get("Evaluation Date"), errors="coerce"
        ).dt.strftime("%d-%b")
        heat = pd.pivot_table(
            fail_rows,
            index="Parameter",
            columns="Date Label",
            values="Result",
            aggfunc="count",
            fill_value=0,
        )
    row_count = len(heat.index) if not heat.empty else 0
    fig_height = max(4, row_count * 0.5)
    fig_heat, axh = plt.subplots(figsize=(6, fig_height))
    if heat.empty:
        axh.text(0.5, 0.5, "No failures recorded", ha="center", va="center", color=theme["text"])
        axh.axis("off")
    else:
        im = axh.imshow(heat.values, cmap="YlOrRd", aspect="auto")
        axh.set_xticks(range(len(heat.columns)))
        axh.set_xticklabels(heat.columns, color=theme["text"])
        axh.set_yticks(range(len(heat.index)))
        axh.set_yticklabels(heat.index, fontsize=8, color=theme["text"])
        for i in range(len(heat.index)):
            for j in range(len(heat.columns)):
                axh.text(
                    j,
                    i,
                    str(int(heat.iloc[i, j])),
                    ha="center",
                    va="center",
                    color=theme["text"],
                    fontsize=8,
                )
        axh.set_title("Failure Distribution by Parameter", fontweight="bold", fontsize=11)
        cbar = plt.colorbar(im, ax=axh)
        cbar.ax.yaxis.set_tick_params(color=theme["text"])
        plt.setp(cbar.ax.get_yticklabels(), color=theme["text"])
        style_chart(axh, theme)
    fig_heat.patch.set_facecolor(theme["bg"])
    plt.tight_layout()
    # 3. Pass vs Fail Pie Chart
    pie_figsize = (6, 6)
    pass_points = summary["Passed Points"].sum() if "Passed Points" in summary.columns else 0
    fail_points = summary["Failed Points"].sum() if "Failed Points" in summary.columns else 0
    fig_pie, axp = plt.subplots(figsize=pie_figsize)
    pass_fail_colors = [theme["pass"], theme["fail"]]
    axp.pie(
        [pass_points, fail_points],
        labels=["Pass", "Fail"],
        autopct="%1.1f%%",
        colors=pass_fail_colors,
        startangle=90,
        radius=0.78,
        labeldistance=1.30,
        pctdistance=1.12,
        wedgeprops={"edgecolor": theme["border"], "linewidth": 1.5},
        textprops={"color": theme["text"], "fontsize": 11, "weight": "bold"},
    )
    axp.set_title("Pass vs Fail Points", fontweight="bold", fontsize=12)
    axp.set_aspect("equal")
    axp.grid(False)
    style_chart(axp, theme)
    fig_pie.patch.set_facecolor(theme["bg"])
    plt.tight_layout()
    # 4. QA Average Scores
    fig_qa, axq = plt.subplots(figsize=(7, 4.5))
    qa_scores = summary.groupby("QA Name")["Overall Score %"].mean().sort_values(ascending=False)
    axq.bar(qa_scores.index, qa_scores.values, color=theme["primary"])
    axq.set_title("QA Average Score (%)", fontweight="bold", fontsize=11)
    axq.set_ylabel("Score %")
    axq.tick_params(axis="x", rotation=20)
    add_bar_labels(axq)
    style_chart(axq, theme)
    fig_qa.patch.set_facecolor(theme["bg"])
    plt.tight_layout()
    # 5. Score per Date
    fig_score_date, axsd = plt.subplots(figsize=(7, 4.5))
    score_by_date = summary.groupby(summary["Evaluation Date"].dt.date)["Overall Score %"].mean()
    score_date_labels = [pd.to_datetime(d).strftime("%d-%b") for d in score_by_date.index]
    axsd.bar(score_date_labels, score_by_date.values, color=theme["accent"])
    axsd.set_title("Average Score by Date (%)", fontweight="bold", fontsize=11)
    axsd.tick_params(axis="x", rotation=30)
    add_bar_labels(axsd)
    style_chart(axsd, theme)
    fig_score_date.patch.set_facecolor(theme["bg"])
    plt.tight_layout()
    # 6. Score per Month
    fig_score_month, axsm = plt.subplots(figsize=(7, 4.5))
    score_by_month = (
        summary
        .assign(Month=pd.to_datetime(summary["Evaluation Date"], errors="coerce").dt.to_period("M"))
        .groupby("Month")["Overall Score %"]
        .mean()
        .sort_index()
    )
    score_month_labels = [
        m.strftime("%b-%y") if pd.notna(m) else ""
        for m in score_by_month.index.to_timestamp()
    ]
    axsm.bar(score_month_labels, score_by_month.values, color=theme["primary"])
    axsm.set_title("Average Score by Month (%)", fontweight="bold", fontsize=11)
    axsm.tick_params(axis="x", rotation=20)
    add_bar_labels(axsm)
    style_chart(axsm, theme)
    fig_score_month.patch.set_facecolor(theme["bg"])
    plt.tight_layout()
    # 7. Audits per Date
    fig_audit_date, axad = plt.subplots(figsize=(7, 4.5))
    audits_by_date = summary.groupby(summary["Evaluation Date"].dt.date).size()
    audit_date_labels = [pd.to_datetime(d).strftime("%d-%b") for d in audits_by_date.index]
    axad.bar(audit_date_labels, audits_by_date.values, color=theme["accent"])
    axad.set_title("Audits per Date", fontweight="bold", fontsize=11)
    axad.tick_params(axis="x", rotation=30)
    add_bar_labels(axad)
    style_chart(axad, theme)
    fig_audit_date.patch.set_facecolor(theme["bg"])
    plt.tight_layout()
    # 8. Audits per Month
    fig_audit_month, axam = plt.subplots(figsize=(7, 4.5))
    audits_by_month = (
        summary
        .assign(Month=pd.to_datetime(summary["Evaluation Date"], errors="coerce").dt.to_period("M"))
        .groupby("Month")
        .size()
        .sort_index()
    )
    audit_month_labels = [
        m.strftime("%b-%y") if pd.notna(m) else ""
        for m in audits_by_month.index.to_timestamp()
    ]
    axam.bar(audit_month_labels, audits_by_month.values, color=theme["primary"])
    axam.set_title("Audits per Month", fontweight="bold", fontsize=11)
    axam.tick_params(axis="x", rotation=20)
    add_bar_labels(axam)
    style_chart(axam, theme)
    fig_audit_month.patch.set_facecolor(theme["bg"])
    plt.tight_layout()
    # 9. Most Failed Parameters
    failed_params = (
        details[details["Result"] == "Fail"]["Parameter"]
        .value_counts()
        .head(10)
        .sort_values(ascending=True)
    )
    bar_count = len(failed_params)
    fig_height = max(4, bar_count * 0.6)
    fig_failed, axf = plt.subplots(figsize=(6, fig_height))
    if failed_params.empty:
        axf.text(0.5, 0.5, "No failures recorded", ha="center", va="center", color=theme["text"])
        axf.axis("off")
    else:
        axf.barh(failed_params.index, failed_params.values, color=theme["fail"])
        axf.set_title("Most Failed Parameters (Top 10)", fontweight="bold", fontsize=11)
        for i, value in enumerate(failed_params.values):
            axf.text(value + 0.1, i, f"{value}", va="center", fontsize=8, color=theme["text"])
    style_chart(axf, theme)
    fig_failed.patch.set_facecolor(theme["bg"])
    plt.tight_layout()
    # 10. Audits per Disposition (HORIZONTAL BAR)
    disp_counts = summary["Call Disposition"].fillna("Unknown").value_counts().sort_values()
    bar_count = len(disp_counts)
    fig_height = max(4, bar_count * 0.6)
    fig_disp, axd = plt.subplots(figsize=(6, fig_height))

    bars = axd.barh(disp_counts.index, disp_counts.values, color=theme["primary"])

    axd.set_title("Audits per Disposition", fontweight="bold", fontsize=12)

    # Add value labels OUTSIDE bars
    for bar in bars:
        width = bar.get_width()
        axd.text(
            width + max(0.3, width * 0.03),
            bar.get_y() + bar.get_height() / 2,
            f"{int(width)}",
            va="center",
            fontsize=9,
            color=theme["text"],
        )
    style_chart(axd, theme)
    fig_disp.patch.set_facecolor(theme["bg"])
    plt.tight_layout()
    return (
        fig_heat,
        fig_trend,
        fig_pie,
        fig_qa,
        fig_score_date,
        fig_score_month,
        fig_audit_date,
        fig_audit_month,
        fig_failed,
        fig_disp,
        summary,
        details,
    )
def fig_to_png_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150)
    buf.seek(0)
    return buf.read()
def dashboard_pdf(figures, title="ATA Dashboard") -> bytes:
    pdf = PDFReport()
    for idx, fig in enumerate(figures):
        img_bytes = fig_to_png_bytes(fig)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            tmp.write(img_bytes)
            tmp_path = tmp.name
        pdf.add_page()
        pdf.set_font("Arial", "B", 14)
        pdf.cell(0, 10, f"{title} - Chart {idx+1}", ln=True)
        pdf.image(tmp_path, x=10, y=30, w=190)
        os.remove(tmp_path)
    out = pdf.output(dest="S")
    return bytes(out) if isinstance(out, (bytes, bytearray)) else out.encode("latin-1")
def dashboard_ppt(figures, title="ATA Dashboard") -> bytes:
    prs = Presentation()
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    margin_x = Inches(0.35)
    margin_y = Inches(0.35)
    max_w = slide_w - (2 * margin_x)
    max_h = slide_h - (2 * margin_y)

    for fig in figures:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img_bytes = fig_to_png_bytes(fig)

        fig_w, fig_h = fig.get_size_inches()
        fig_aspect = fig_w / fig_h if fig_h else 1.0
        box_aspect = max_w / max_h if max_h else 1.0

        if fig_aspect >= box_aspect:
            pic_w = max_w
            pic_h = int(max_w / fig_aspect)
        else:
            pic_h = max_h
            pic_w = int(max_h * fig_aspect)

        x = int((slide_w - pic_w) / 2)
        y = int((slide_h - pic_h) / 2)
        slide.shapes.add_picture(io.BytesIO(img_bytes), x, y, width=pic_w, height=pic_h)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
# -------------------- MAIN APP --------------------

def clear_login_state(cookie_manager):
    try:
        if cookie_manager.get(COOKIE_AUTH_KEY):
            cookie_manager.delete(COOKIE_AUTH_KEY)
    except:
        pass

    for key in list(st.session_state.keys()):
        st.session_state.pop(key, None)

    st.session_state.authenticated = False

def reset_evaluation_form() -> None:
    st.session_state.edit_mode = False
    st.session_state.edit_eval_id = ""
    st.session_state.prefill = {}
    st.session_state.form_key_id = uuid4().hex
    st.session_state.save_in_progress = False
for key in [
    "edit_mode",
    "edit_eval_id",
    "prefill",
    "goto_nav",
    "last_saved_id",
    "reset_notice",
    "reset_counter",
    "authenticated",
    "theme_mode",
    "form_key_id",
    "save_in_progress",
    "coaching_summary_text",
]:
    if key not in st.session_state:
        if key == "prefill":
            st.session_state[key] = {}
        elif key == "reset_counter":
            st.session_state[key] = 0
        elif key == "authenticated":
            st.session_state[key] = False
        elif key == "theme_mode":
            st.session_state[key] = "system"
        elif key == "form_key_id":
            st.session_state[key] = uuid4().hex
        elif key == "save_in_progress":
            st.session_state[key] = False
        else:
            st.session_state[key] = ""

cookie_manager = stx.CookieManager(key="ata_cookie_manager")
active_theme = get_active_theme()
apply_base_css()
apply_theme_css(active_theme)
secure_authentication_gate()
if st.session_state.get("goto_nav"):
    st.session_state["nav_radio"] = st.session_state["goto_nav"]
    st.session_state["goto_nav"] = ""

st.sidebar.markdown(
    f"""
    <div class="ata-hero">
      <p class="t1">{DAMAC_TITLE}</p>
      <p class="t2">{APP_NAME}</p>
    </div>
    <div style="display:flex;justify-content:center;margin:0.6rem 0 0.2rem 0;">
        <img src="{LOGO_URL}" style="max-width:90%;height:auto;" />
    </div>
    """,
    unsafe_allow_html=True,
)

with st.sidebar.container(border=False):
    st.markdown('<div class="ata-nav-card">', unsafe_allow_html=True)
    theme_choice = st.selectbox(
        "Theme",
        ["System", "Light", "Dark"],
        index={"system": 0, "light": 1, "dark": 2}.get(st.session_state.theme_mode, 0),
        key="theme_selectbox",
    )
    st.session_state.theme_mode = theme_choice.lower()
    active_theme = get_active_theme()
    apply_theme_css(active_theme)

    nav_options = ["Home", "Evaluation", "View", "Dashboard"]
    if st.session_state.get("nav_radio") not in nav_options:
        st.session_state["nav_radio"] = "Home"

    st.markdown('<div class="sidebar-nav-btn">', unsafe_allow_html=True)
    for page in nav_options:
        label = f"● {page}" if st.session_state.get("nav_radio") == page else page
        if st.button(label, key=f"nav_btn_{page}", use_container_width=True):
            st.session_state["nav_radio"] = page
            st.rerun()
    st.markdown('</div>', unsafe_allow_html=True)

    if st.button("🚪 Logout", use_container_width=True):
        clear_login_state(cookie_manager)
        st.rerun()
    st.markdown('</div>', unsafe_allow_html=True)

st.sidebar.markdown(
    '<div style="text-align:center;color:var(--accent-gold);font-size:13px;margin-top:0.9rem;">Designed and built by Mohamed Seddiq</div>',
    unsafe_allow_html=True,
)

nav = st.session_state.get("nav_radio", "Home")
if nav == "Home":
    summary = read_google_summary()
    render_title_card("Home")
    st.markdown(
        f"""
        <div class="ata-hero left-align">
          <p class="t1">Welcome to {APP_NAME}</p>
          <p class="t2">Monitor, evaluate, and improve auditor performance with real-time insights.</p>
        </div>
        """,
        unsafe_allow_html=True,
    )
    c1, c2, c3, c4 = st.columns(4)
    if summary.empty:
        for c, l in zip([c1, c2, c3, c4], ["Total Evals", "Avg Score", "Failure Rate", "Last Audit"]):
            with c:
                st.markdown(
                    f"""
                    <div class="stat-card">
                      <div class="stat-val">0</div>
                      <div class="stat-label">{l}</div>
                    </div>
                    """,
                    unsafe_allow_html=True,
                )
    else:
        for col in ["Failed Points", "Total Points", "Overall Score %"]:
            summary[col] = pd.to_numeric(summary[col], errors="coerce").fillna(0)
        stats = [
            len(summary),
            f"{summary['Overall Score %'].mean():.2f}%",
            (
                f"{(summary['Failed Points'].sum() / summary['Total Points'].sum() * 100 if summary['Total Points'].sum() else 0):.2f}%"
            ),
            format_date(summary.iloc[-1]["Evaluation Date"]),
        ]
        labels = ["Total Evaluations", "Average Score", "Avg Failure Rate", "Last Evaluation"]
        for c, v, l in zip([c1, c2, c3, c4], stats, labels):
            with c:
                st.markdown(
                    f"""
                    <div class="stat-card">
                      <div class="stat-val">{v}</div>
                      <div class="stat-label">{l}</div>
                    </div>
                    """,
                    unsafe_allow_html=True,
                )
    st.write("")
    st.write("")
    b1, b2, b3 = st.columns(3)
    with b1:
        st.markdown('<div class="action-card">', unsafe_allow_html=True)
        if st.button("➕ New Evaluation", use_container_width=True):
            st.session_state.goto_nav = "Evaluation"
            st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)
    with b2:
        st.markdown('<div class="action-card">', unsafe_allow_html=True)
        if st.button("🔍 View Records", use_container_width=True):
            st.session_state.goto_nav = "View"
            st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)
    with b3:
        st.markdown('<div class="action-card">', unsafe_allow_html=True)
        if st.button("📊 Performance Dashboard", use_container_width=True):
            st.session_state.goto_nav = "Dashboard"
            st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)
    if not summary.empty:
        st.markdown("### Recent Audited Transactions")
        recent = summary.sort_values("Evaluation Date", ascending=False).head(3).copy()
        recent["Evaluation Date"] = recent["Evaluation Date"].apply(format_date)
        recent["Audit Date"] = recent["Audit Date"].apply(format_date)
        recent["Overall Score %"] = (
            pd.to_numeric(recent["Overall Score %"], errors="coerce")
            .fillna(0)
            .map("{:.2f}%".format)
        )
        st.dataframe(
            recent[
                [
                    "Evaluation ID",
                    "Evaluation Date",
                    "QA Name",
                    "Auditor",
                    "Audit Date",
                    "Call ID",
                    "Overall Score %",
                ]
            ],
            use_container_width=True,
            hide_index=True,
        )
elif nav == "Evaluation":
    render_title_card(("Edit" if st.session_state.edit_mode else "New") + " Evaluation")
    if st.session_state.get("reset_notice"):
        st.success(st.session_state.reset_notice)
        st.session_state.reset_notice = ""
    pre, df_all = st.session_state.prefill or {}, load_parameters_df()
    if st.session_state.edit_mode and isinstance(pre.get("details_df"), pd.DataFrame):
        df_all = pre["details_df"].copy()
    df_all = normalize_details_df(df_all)
    df_all["Comment"] = df_all.get("Comment", "").fillna("")
    form_token = st.session_state.form_key_id
    with st.form(f"eval_form_{form_token}"):
        c1, c2, c3, c4 = st.columns(4)
        qa_name = c1.text_input("QA Name", value=pre.get("qa_name", ""), key=f"qa_name_{form_token}")
        auditor = c1.text_input("Auditor", value=pre.get("auditor", ""), key=f"auditor_{form_token}")
        eval_date = c2.date_input(
            "Eval Date",
            value=pre.get("evaluation_date", date.today()),
            format="DD/MM/YYYY",
            key=f"eval_date_{form_token}",
        )
        audit_date = c2.date_input(
            "Audit Date",
            value=pre.get("audit_date", date.today()),
            format="DD/MM/YYYY",
            key=f"audit_date_{form_token}",
        )
        call_id = c3.text_input("Call ID", value=pre.get("call_id", ""), key=f"call_id_{form_token}")
        call_dur = c3.text_input("Duration", value=pre.get("call_duration", ""), key=f"call_duration_{form_token}")
        call_disp = c4.text_input("Disposition", value=pre.get("call_disposition", ""), key=f"call_disposition_{form_token}")
        reaudit = c4.selectbox(
            "Reaudit",
            ["No", "Yes"],
            index=0 if pre.get("reaudit") != "Yes" else 1,
            key=f"reaudit_{form_token}",
        )
        df_acc = df_all[df_all["Group"] == "ACCURACY_SUB"].copy()
        df_qual = df_all[df_all["Group"] == "EVAL_QUALITY"].copy()
        editor_key = f"ed_{form_token}"
        with st.expander("Accuracy of Scoring", expanded=True):
            ed_acc = st.data_editor(
                df_acc[["Parameter", "Result", "Comment"]],
                use_container_width=True,
                key=f"{editor_key}_acc",
                column_config={
                    "Result": st.column_config.SelectboxColumn(
                        options=["Pass", "Fail"],
                        required=True,
                        help="Select Pass or Fail.",
                    ),
                    "Comment": st.column_config.TextColumn(),
                },
            )
        with st.expander("Evaluation Quality", expanded=True):
            ed_qual = st.data_editor(
                df_qual[["Parameter", "Result", "Comment"]],
                use_container_width=True,
                key=f"{editor_key}_qual",
                column_config={
                    "Result": st.column_config.SelectboxColumn(
                        options=["Pass", "Fail"],
                        required=True,
                        help="Select Pass or Fail.",
                    ),
                    "Comment": st.column_config.TextColumn(),
                },
            )
        st.markdown('<div class="eval-action-row">', unsafe_allow_html=True)
        action_cols = st.columns(3 if st.session_state.edit_mode else 2)
        with action_cols[0]:
            save_clicked = st.form_submit_button("💾 Save Evaluation", use_container_width=True, disabled=st.session_state.get("save_in_progress", False))
        with action_cols[1]:
            reset_clicked = st.form_submit_button("🔄 Reset Form", use_container_width=True)
        cancel_clicked = False
        if st.session_state.edit_mode:
            with action_cols[2]:
                cancel_clicked = st.form_submit_button("↩️ Cancel Edit", use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)
        if reset_clicked:
            reset_evaluation_form()
            st.session_state.reset_notice = "Form reset."
            st.rerun()
        if cancel_clicked:
            reset_evaluation_form()
            st.session_state.goto_nav = "View"
            st.rerun()
        if save_clicked:
            if st.session_state.get("save_in_progress", False):
                st.info("Save already in progress. Please wait.")
                st.stop()
            if "Parameter" not in ed_acc.columns or "Parameter" not in ed_qual.columns:
                st.error("Unable to save. Please reset the form and try again.")
                st.stop()
            eval_id = (
                st.session_state.edit_eval_id
                if st.session_state.edit_mode
                else next_evaluation_id(eval_date.strftime("%Y-%m-%d"))
            )
            st.session_state.save_in_progress = True
            try:
                acc_full, qual_full = df_acc.copy(), df_qual.copy()
                acc_full["Result"], acc_full["Comment"] = ed_acc["Result"].values, ed_acc["Comment"].values
                qual_full["Result"], qual_full["Comment"] = (
                    ed_qual["Result"].values,
                    ed_qual["Comment"].values,
                )
                details_all = pd.concat([acc_full, qual_full], ignore_index=True)
                metrics = compute_weighted_score(details_all)
                record = {
                    "evaluation_id": eval_id,
                    "qa_name": qa_name,
                    "auditor": auditor,
                    "evaluation_date": eval_date.strftime("%Y-%m-%d"),
                    "audit_date": audit_date.strftime("%Y-%m-%d"),
                    "reaudit": reaudit,
                    "call_id": call_id,
                    "call_duration": call_dur,
                    "call_disposition": call_disp,
                    "overall_score": metrics["score"],
                    "passed_points": metrics["passed_points"],
                    "failed_points": metrics["failed_points"],
                    "total_points": metrics["total_points"],
                    "details": details_all,
                }
                upsert_google_sheet(record)
                was_edit_mode = st.session_state.edit_mode
                reset_evaluation_form()
                st.session_state.last_saved_id = eval_id
                if was_edit_mode:
                    st.session_state.goto_nav = "View"
                    st.session_state.reset_notice = "Evaluation updated successfully."
                else:
                    st.session_state.reset_notice = f"Saved Evaluation ID: {eval_id}. Form reset for a new entry."
                st.rerun()
            finally:
                st.session_state.save_in_progress = False
elif nav == "View":
    render_title_card("Audit Records Explorer")
    summary = read_google_summary()
    details = read_google_details()
    if summary.empty:
        st.info("No records found.")
    else:
        if st.session_state.get("last_saved_id"):
            st.success(f"Saved Evaluation ID: {st.session_state.last_saved_id}")
            st.session_state.last_saved_id = ""
        summary["QA Name"] = summary["QA Name"].fillna("N/A")
        summary["Audit Date"] = summary["Audit Date"].fillna("N/A")
        qas = sorted(summary["QA Name"].unique().tolist())
        dates = sorted([str(d) for d in summary["Audit Date"].unique().tolist()])
        f1, f2, f3 = st.columns(3)
        qa_f = f1.selectbox("Filter by QA", ["All"] + qas)
        dt_f = f2.selectbox("Filter by Date", ["All"] + dates)
        search = f3.text_input("Search (ID/Auditor/Call)")
        filtered = summary.copy()
        if qa_f != "All":
            filtered = filtered[filtered["QA Name"] == qa_f]
        if dt_f != "All":
            filtered = filtered[filtered["Audit Date"].astype(str) == dt_f]
        if search:
            filtered = filtered[filtered.apply(lambda r: search.lower() in str(r).lower(), axis=1)]
        if filtered.empty:
            st.warning("No matches.")
        else:
            display_df = filtered.copy()
            display_df.insert(0, "S.No", range(1, len(display_df) + 1))
            display_df["Evaluation Date"] = display_df["Evaluation Date"].apply(format_date)
            display_df["Audit Date"] = display_df["Audit Date"].apply(format_date)
            display_df["Overall Score %"] = (
                pd.to_numeric(display_df["Overall Score %"], errors="coerce")
                .fillna(0)
                .map("{:.2f}%".format)
            )
            summary_cols = [
                "S.No",
                "Evaluation ID",
                "Evaluation Date",
                "QA Name",
                "Auditor",
                "Audit Date",
                "Call ID",
                "Overall Score %",
            ]
            if active_theme["mode"] == "dark":
                summary_html = display_df[summary_cols].to_html(index=False, classes="summary-records-table", border=0)
                st.markdown(summary_html, unsafe_allow_html=True)
            else:
                st.dataframe(
                    display_df[summary_cols],
                    use_container_width=True,
                    hide_index=True,
                )
            record_options = display_df.apply(
                lambda r: f"{r['S.No']} | {r['Evaluation ID']} | {r['QA Name']}", axis=1
            ).tolist()
            sel_label = st.selectbox("Select Record to View Details", record_options)
            sel_id = sel_label.split(" | ")[1]
            if st.session_state.get("coaching_summary_eval_id") != str(sel_id).strip():
                st.session_state.coaching_summary_text = ""
                st.session_state.coaching_summary_eval_id = str(sel_id).strip()
            selected_rows = summary[summary["Evaluation ID"].astype(str).str.strip() == str(sel_id).strip()]
            if selected_rows.empty:
                st.error("Selected record was not found. Please refresh and try again.")
                st.stop()
            row = selected_rows.iloc[0]
            eval_date_display = format_date(row["Evaluation Date"])
            audit_date_display = format_date(row["Audit Date"])
            email_subject = f"ATA Evaluation | {sel_id} | {row['QA Name']} | {audit_date_display}"
            st.markdown(
                f"""
                <div class="ata-card">
                  <h3 class="view-detail-title" style="margin-top:0;">Evaluation Details: {sel_id}</h3>
                  <table class="styled-table eval-details-table">
                    <tr><td><b>Evaluation ID:</b> {sel_id}</td><td><b>Evaluation Date:</b> {eval_date_display}</td></tr>
                    <tr><td><b>QA Name:</b> {row['QA Name']}</td><td><b>Auditor Name:</b> {row['Auditor']}</td></tr>
                    <tr><td><b>Audit Date:</b> {audit_date_display}</td><td><b>Call ID:</b> {row['Call ID']}</td></tr>
                    <tr><td><b>Call Duration:</b> {row['Call Duration']}</td><td><b>Call Disposition:</b> {row['Call Disposition']}</td></tr>
                    <tr><td><b>Overall Score:</b> <span class="view-score" style="font-size:18px;font-weight:bold;">{row['Overall Score %']:.2f}%</span></td><td><b>Reaudit:</b> <span style="font-weight:700;color:{'#d93025' if str(row['Reaudit']).strip().lower() == 'yes' else '#1f8f4a'};">{row['Reaudit']}</span></td></tr>
                    <tr><td colspan="2"><b>Email Subject:</b> {email_subject}</td></tr>
                  </table>
                </div>
                """,
                unsafe_allow_html=True,
            )
            selected_details = details[details["Evaluation ID"].astype(str).str.strip() == str(sel_id).strip()].copy()
            detail_cols = ["Group", "Parameter", "Result", "Comment", "Description"]
            available_detail_cols = [c for c in detail_cols if c in selected_details.columns]
            if available_detail_cols:
                selected_details = selected_details[available_detail_cols]
            rec = {
                "evaluation_id": sel_id,
                "qa_name": row["QA Name"],
                "auditor": row["Auditor"],
                "evaluation_date": eval_date_display,
                "audit_date": audit_date_display,
                "reaudit": row["Reaudit"],
                "call_id": row["Call ID"],
                "call_duration": row["Call Duration"],
                "call_disposition": row["Call Disposition"],
                "overall_score": row["Overall Score %"],
                "details": selected_details,
            }
            export_buf = io.BytesIO()
            with pd.ExcelWriter(export_buf, engine="openpyxl") as writer:
                pd.DataFrame([row]).to_excel(writer, sheet_name="Summary", index=False)
                selected_details.to_excel(
                    writer, sheet_name="Details", index=False
                )
            export_buf.seek(0)

            st.markdown(
                """
                <style>
                .view-action-grid {
                    display: grid;
                    grid-template-columns: repeat(3, 1fr);
                    grid-template-rows: repeat(3, auto);
                    gap: 12px;
                    width: 100%;
                }
                .view-action-grid .stButton > button,
                .view-action-grid .stDownloadButton > button,
                .view-action-grid .stForm [data-testid="stFormSubmitButton"] > button {
                    height: 60px !important;
                    min-height: 60px !important;
                    max-height: 60px !important;
                    border-radius: 10px !important;
                    font-size: 14px !important;
                    font-weight: 600 !important;
                    display: flex !important;
                    align-items: center !important;
                    justify-content: center !important;
                    white-space: nowrap !important;
                    overflow: hidden !important;
                    text-overflow: ellipsis !important;
                }
                .view-action-grid iframe {
                    width: 100% !important;
                    height: 60px !important;
                    border: 0 !important;
                }
                </style>
                """,
                unsafe_allow_html=True,
            )

            st.markdown('<div class="view-action-grid">', unsafe_allow_html=True)

            row1 = st.columns(3, gap="small")
            with row1[0]:
                pdf_clicked = st.download_button(
                    "📄 Download PDF",
                    pdf_evaluation(rec),
                    f"ATA_{sel_id}.pdf",
                    "application/pdf",
                    use_container_width=True,
                )
                if pdf_clicked:
                    st.success("PDF Ready")
            with row1[1]:
                copy_html_to_clipboard_button("📋 Copy Email Body", email_html_inline(rec), f"copy_body_{sel_id}", active_theme)
            with row1[2]:
                copy_html_to_clipboard_button("📌 Copy Email Subject", email_subject_text(rec), f"copy_subject_{sel_id}", active_theme)

            row2 = st.columns(3, gap="small")
            with row2[0]:
                if st.button("✏️ Edit Record", use_container_width=True):
                    st.session_state.edit_mode = True
                    st.session_state.edit_eval_id = sel_id
                    st.session_state.prefill = {
                        "qa_name": row["QA Name"],
                        "auditor": row["Auditor"],
                        "evaluation_date": pd.to_datetime(row["Evaluation Date"]).date(),
                        "audit_date": pd.to_datetime(row["Audit Date"]).date(),
                        "call_id": row["Call ID"],
                        "call_duration": row["Call Duration"],
                        "call_disposition": row["Call Disposition"],
                        "reaudit": row["Reaudit"],
                        "details_df": details[details["Evaluation ID"].astype(str).str.strip() == str(sel_id).strip()],
                    }
                    st.session_state.goto_nav = "Evaluation"
                    st.rerun()
            with row2[1]:
                if st.button("🗑️ Delete Record", use_container_width=True):
                    if delete_evaluation(sel_id):
                        st.success(f"Deleted {sel_id}")
                        st.rerun()
            with row2[2]:
                export_clicked = st.download_button(
                    "📥 Export Selected to Excel",
                    export_buf,
                    f"ATA_{sel_id}.xlsx",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True,
                )
                if export_clicked:
                    st.success("Excel Ready")

            row3 = st.columns(3, gap="small")
            with row3[0]:
                if st.button("🧠 Generate Coaching Summary", use_container_width=True):
                    auditor_base = compute_auditor_intelligence(summary, details)
                    auditor_risk = compute_risk_flags(auditor_base, details)
                    row_metrics = auditor_risk[auditor_risk["Auditor"].astype(str).str.strip() == str(row["Auditor"]).strip()]
                    metric_payload = row_metrics.iloc[0].to_dict() if not row_metrics.empty else {"Risk Level": "Low"}
                    st.session_state.coaching_summary_text = generate_coaching_summary(rec, metric_payload)
                    st.session_state.coaching_summary_eval_id = str(sel_id).strip()
                    st.success("Coaching Summary Generated")
            with row3[1]:
                if st.session_state.get("coaching_summary_text"):
                    copy_html_to_clipboard_button(
                        "📋 Copy Coaching Summary",
                        f"<pre>{st.session_state.get('coaching_summary_text', '')}</pre>",
                        f"copy_coach_{sel_id}",
                        active_theme,
                    )
                else:
                    st.button("📋 Copy Coaching Summary", use_container_width=True, disabled=True)
            with row3[2]:
                if st.button("🧹 Clear Coaching Summary", use_container_width=True, disabled=not st.session_state.get("coaching_summary_text")):
                    st.session_state.coaching_summary_text = ""
                    st.warning("Coaching Summary Cleared")
                    st.rerun()

            st.markdown('</div>', unsafe_allow_html=True)

            if st.session_state.get("coaching_summary_text"):
                st.text_area(
                    "Coaching Summary",
                    value=st.session_state.get("coaching_summary_text", ""),
                    height=260,
                    key=f"coach_text_{sel_id}",
                )

            st.markdown("<div class='group-title'>Parameter Breakdown</div>", unsafe_allow_html=True)
            det = details[details["Evaluation ID"].astype(str).str.strip() == str(sel_id).strip()]
            if not det.empty:
                for grp in ["ACCURACY_SUB", "EVAL_QUALITY"]:
                    with st.expander(grp.replace("_", " ").title(), expanded=True):
                        grp_df = det[det["Group"] == grp][["Parameter", "Result", "Comment"]].copy()
                        if active_theme["mode"] == "dark":
                            st.markdown(grp_df.to_html(index=False, classes="dark-breakdown-table", border=0), unsafe_allow_html=True)
                        else:
                            st.table(grp_df)
elif nav == "Dashboard":
    render_title_card("Performance Dashboard", "Visualizing quality trends and failure distributions.")
    summary_all = read_google_summary()
    details_all = read_google_details()
    if not summary_all.empty and "Evaluation Date" in summary_all.columns:
        summary_all["Evaluation Date"] = pd.to_datetime(summary_all["Evaluation Date"], errors="coerce")
    if not details_all.empty and "Evaluation Date" in details_all.columns:
        details_all["Evaluation Date"] = pd.to_datetime(details_all["Evaluation Date"], errors="coerce")
    if summary_all.empty or details_all.empty:
        st.info("No data available for analysis.")
    else:
        filter_col1, filter_col2, filter_col3, filter_col4 = st.columns(4)
        qa_options = ["All"] + sorted(summary_all["QA Name"].dropna().unique().tolist())
        disp_options = ["All"] + sorted(summary_all["Call Disposition"].dropna().unique().tolist())
        month_options = ["All"] + sorted(summary_all["Evaluation Date"].dt.to_period("M").astype(str).unique().tolist())
        date_options = ["All"] + sorted(summary_all["Evaluation Date"].dt.strftime("%d-%b").dropna().unique().tolist())
        qa_filter = filter_col1.selectbox("Filter by QA", qa_options)
        disp_filter = filter_col2.selectbox("Filter by Disposition", disp_options)
        month_filter = filter_col3.selectbox("Filter by Month", month_options)
        date_filter = filter_col4.selectbox("Filter by Date", date_options)
        summary = summary_all.copy()
        if qa_filter != "All":
            summary = summary[summary["QA Name"] == qa_filter]
        if disp_filter != "All":
            summary = summary[summary["Call Disposition"] == disp_filter]
        if month_filter != "All":
            summary = summary[summary["Evaluation Date"].dt.to_period("M").astype(str) == month_filter]
        if date_filter != "All":
            summary = summary[summary["Evaluation Date"].dt.strftime("%d-%b") == date_filter]
        details = details_all.copy()
        if "Evaluation ID" in summary.columns:
            details = details[
                details["Evaluation ID"].astype(str).str.strip().isin(
                    summary["Evaluation ID"].astype(str).str.strip().unique()
                )
            ]
        if summary.empty or details.empty:
            st.info("No data available for analysis.")
        else:
            (
                fig_heat,
                fig_trend,
                fig_pie,
                fig_qa,
                fig_score_date,
                fig_score_month,
                fig_audit_date,
                fig_audit_month,
                fig_failed,
                fig_disp,
                _summary_unused,
                _details_unused,
            ) = build_dashboard_figs(summary, details)
            if fig_trend:
                rows = [
                    (fig_pie, fig_disp),
                    (fig_trend, fig_qa),
                    (fig_score_month, fig_score_date),
                    (fig_audit_month, fig_audit_date),
                    (fig_heat, fig_failed),
                ]

                for left_fig, right_fig in rows:
                    col1, col2 = st.columns(2, gap="large")
                    with col1:
                        st.pyplot(left_fig, use_container_width=False)
                    with col2:
                        st.pyplot(right_fig, use_container_width=False)

                st.markdown("### Auditor Performance Intelligence")
                auditor_intel = compute_auditor_intelligence(summary, details)
                risk_df = compute_risk_flags(auditor_intel, details)
                health_df = compute_health_index(auditor_intel, details)
                ranking_df = auditor_intel.merge(health_df[["Auditor", "Health Index", "Health Classification"]], on="Auditor", how="left")
                ranking_df["Health Index"] = pd.to_numeric(
                    ranking_df["Health Index"],
                    errors="coerce"
                ).fillna(0)
                ranking_df = ranking_df.sort_values("Health Index", ascending=False)
                ranking_display = ranking_df.copy()
                if "Volatility" in ranking_display.columns:
                    ranking_display = ranking_display.drop(columns=["Volatility"])
                for pct_col in ["Avg Score", "Failure Rate", "Reaudit Ratio"]:
                    if pct_col in ranking_display.columns:
                        vals = pd.to_numeric(ranking_display[pct_col], errors="coerce").fillna(0)
                        ranking_display[pct_col] = vals.round(0).astype(int).astype(str) + "%"
                st.dataframe(ranking_display, use_container_width=True, hide_index=True)

                st.markdown("### 🚩 Interactions Requiring Coaching")
                interactions = summary.copy()
                interactions["Overall Score %"] = pd.to_numeric(interactions.get("Overall Score %", 0), errors="coerce").fillna(0)
                interactions["Evaluation Date"] = pd.to_datetime(interactions.get("Evaluation Date"), errors="coerce")
                detail_fail = details[details.get("Result", "").astype(str).str.lower() == "fail"].copy()
                if not detail_fail.empty:
                    detail_fail["Parameter"] = detail_fail.get("Parameter", "").fillna("").astype(str).str.strip()
                    failed_params = detail_fail.groupby("Evaluation ID")["Parameter"].apply(
                        lambda x: " | ".join(sorted(dict.fromkeys([p for p in x if p]))) if any([p for p in x if p]) else "No Failed Parameters"
                    ).reset_index(name="Failed Parameters")
                else:
                    failed_params = pd.DataFrame(columns=["Evaluation ID", "Failed Parameters"])

                interactions = interactions.merge(failed_params, on="Evaluation ID", how="left")
                if "Reaudit" not in interactions.columns:
                    interactions["Reaudit"] = ""
                critical_ids = details[
                    details.get("Parameter", "").astype(str).str.strip().str.lower().eq("critical error identification")
                    & details.get("Result", "").astype(str).str.strip().str.lower().eq("fail")
                ]["Evaluation ID"].astype(str).str.strip().unique().tolist()
                repeat_by_auditor = auditor_intel.set_index("Auditor")["Repeat Failure Count"] if not auditor_intel.empty else pd.Series(dtype=float)
                interactions["_repeat"] = interactions.get("Auditor", "").map(repeat_by_auditor).fillna(0)
                interactions["_critical"] = interactions["Evaluation ID"].astype(str).str.strip().isin(critical_ids)
                interactions["_score"] = interactions["Overall Score %"] < 85
                interactions["_reaudit"] = interactions["Reaudit"].astype(str).str.strip().str.lower().eq("yes")
                interactions = interactions[(interactions["_score"]) | (interactions["_critical"]) | (interactions["_repeat"] >= 1) | (interactions["_reaudit"])]

                def _reason(r):
                    reasons = []
                    if r["_score"]:
                        reasons.append("Score")
                    if r["_critical"]:
                        reasons.append("Critical")
                    if r["_repeat"] >= 1:
                        reasons.append("Repeat")
                    if r["_reaudit"]:
                        reasons.append("Reaudit")
                    return " / ".join(reasons)

                if interactions.empty:
                    st.info("No high-risk interactions detected.")
                else:
                    interactions["Trigger Reason"] = interactions.apply(_reason, axis=1)
                    qa_map = risk_df.set_index("Auditor")["QA Intervention Required"] if not risk_df.empty else pd.Series(dtype=bool)
                    coach_map = risk_df.set_index("Auditor")["Coaching Required"] if not risk_df.empty else pd.Series(dtype=bool)
                    interactions["QA Intervention Required"] = interactions.get("Auditor", "").map(qa_map).fillna(False).map(lambda x: "Yes" if bool(x) else "No")
                    interactions["Coaching Required"] = interactions.get("Auditor", "").map(coach_map).fillna(False).map(lambda x: "Yes" if bool(x) else "No")

                    interactions = interactions.sort_values("Evaluation Date", ascending=False)

                    out_cols = [
                        "Evaluation ID",
                        "Evaluation Date",
                        "QA Name",
                        "Auditor",
                        "Overall Score %",
                        "Failed Parameters",
                        "Coaching Required",
                        "QA Intervention Required",
                        "Trigger Reason",
                    ]
                    for c in out_cols:
                        if c not in interactions.columns:
                            interactions[c] = ""
                    st.dataframe(interactions[out_cols], use_container_width=True, hide_index=True)

                    st.markdown(
                        """
                        **Coaching Guideline Example**
                        - Conduct structured 1:1 feedback.
                        - Focus on failed parameters.
                        - Reinforce QA standard alignment.
                        - Monitor next 3 evaluations.

                        **QA Intervention Guideline Example**
                        - Conduct recalibration review.
                        - Validate scoring consistency.
                        - Review previous 5 audits.
                        - Escalate if pattern continues.
                        """
                    )
                st.divider()
                st.markdown('<div class="dashboard-action-btn">', unsafe_allow_html=True)
                d1, d2, d3 = st.columns(3)
                figures = [
                    fig_pie,
                    fig_disp,
                    fig_trend,
                    fig_qa,
                    fig_score_month,
                    fig_score_date,
                    fig_audit_month,
                    fig_audit_date,
                    fig_heat,
                    fig_failed,
                ]
                with d1:
                    st.download_button(
                        "📥 Download Dashboard PDF",
                        dashboard_pdf(figures),
                        "ATA_Dashboard.pdf",
                        "application/pdf",
                        use_container_width=True,
                    )
                with d2:
                    st.download_button(
                        "📥 Download Dashboard PPT",
                        dashboard_ppt(figures),
                        "ATA_Dashboard.pptx",
                        use_container_width=True,
                    )
                with d3:
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp_file:
                        temp_xlsx = tmp_file.name
                    with pd.ExcelWriter(temp_xlsx, engine="openpyxl") as writer:
                        summary_all.to_excel(writer, sheet_name="Summary", index=False)
                        details_all.to_excel(writer, sheet_name="Details", index=False)
                    write_formatted_report({"evaluation_id": "", "details": pd.DataFrame()}, temp_xlsx, summary_all, details_all)
                    with open(temp_xlsx, "rb") as f:
                        excel_bytes = f.read()
                    os.remove(temp_xlsx)
                    st.download_button(
                        "📥 Download Excel Log",
                        excel_bytes,
                        "ATA_Audit_Log.xlsx",
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        use_container_width=True,
                    )
                st.markdown('</div>', unsafe_allow_html=True)























